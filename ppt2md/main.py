"""CLI entry point for ppt2md."""

import argparse
import os
import sys
from pathlib import Path

from ppt2md.parser.presentation import open_presentation, get_slide_count
from ppt2md.parser.text import extract_text_from_slide
from ppt2md.parser.placeholder import get_placeholder_type, placeholder_to_markdown_prefix
from ppt2md.parser.image import extract_images_from_slide
from ppt2md.parser.table import extract_tables_from_slide
from ppt2md.parser.chart import extract_charts_from_slide
from ppt2md.parser.chart_export import export_charts_from_slide
from ppt2md.parser.formula import find_omml_elements, convert_omml_to_latex
from ppt2md.parser.notes import extract_notes_from_slide, format_notes_markdown
from ppt2md.parser.background import extract_background_image, format_background_markdown
from ppt2md.parser.media import extract_media_from_slide, format_media_markdown, detect_ole_objects, format_ole_markdown
from ppt2md.parser.smartart import extract_group_shapes, extract_group_texts, extract_smartart_text
from ppt2md.parser.sections import get_sections, format_section_markdown

from ppt2md.converter.frontmatter import generate_frontmatter
from ppt2md.converter.sort_utils import sort_shapes_by_position
from ppt2md.converter.filter_utils import is_empty_slide, should_skip_shape
from ppt2md.converter.format_utils import format_paragraph_runs
from ppt2md.converter.list_utils import format_paragraph_as_list
from ppt2md.converter.position_utils import format_image_markdown, format_position_comment
from ppt2md.converter.batch import find_pptx_files, get_output_dir_for_file, batch_convert_summary
from ppt2md.converter.metadata import extract_shape_metadata
from ppt2md.converter.error_handler import (
    handle_file_error, handle_parse_error, handle_image_error,
    handle_formula_error, get_exit_code, EXIT_SUCCESS
)

from pptx.enum.shapes import MSO_SHAPE_TYPE


def convert_pptx_to_markdown(input_path, output_dir=None, include_notes=False,
                              include_hidden=False, skip_empty=False,
                              formula_as_image=False, keep_original_format=False,
                              image_dpi=96, no_frontmatter=False,
                              verbose=False, debug=False):
    """Convert a single PPTX file to Markdown.

    Returns:
        dict with success, output_file, errors, slide_count.
    """
    input_path = Path(input_path)
    result = {"file": str(input_path), "success": False, "errors": [], "slide_count": 0}

    try:
        prs = open_presentation(str(input_path))
    except Exception as e:
        result["errors"].append(str(e))
        handle_file_error(e, input_path.name, verbose)
        return result

    result["slide_count"] = get_slide_count(prs)

    if output_dir is None:
        output_dir = get_output_dir_for_file(input_path)
    else:
        output_dir = Path(output_dir)

    images_dir = output_dir / "images"
    media_dir = output_dir / "media"
    os.makedirs(images_dir, exist_ok=True)

    seen_rids = {}
    md_parts = []
    errors = []

    if not no_frontmatter:
        md_parts.append(generate_frontmatter(prs, input_path.name))

    # Store presentation-level metadata
    import json
    pres_meta = json.dumps({
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
    }, ensure_ascii=False)
    md_parts.append("\n<!-- PPTX_PRESENTATION_META_START\n{}\nPPTX_PRESENTATION_META_END -->".format(pres_meta))

    sections = get_sections(prs)
    section_map = {}
    for sec in sections:
        idx = sec.get("start_slide_idx", 0)
        section_map[idx] = sec

    slide_num = 0
    for slide_idx, slide in enumerate(prs.slides):
        slide_num += 1

        if not include_hidden:
            try:
                if slide.slide_id in [s.slide_id for s in prs.slides._sldIdLst
                                       if s.get("hidden")]:
                    continue
            except (AttributeError, TypeError):
                pass

        if skip_empty and is_empty_slide(slide):
            if verbose:
                print("Skipping empty slide {}".format(slide_num))
            continue

        if slide_idx in section_map:
            md_parts.append("\n---\n\n{}".format(format_section_markdown(section_map[slide_idx])))

        md_parts.append("\n---\n\n## Slide {}".format(slide_num))

        sorted_shapes = sort_shapes_by_position(list(slide.shapes))

        # Build theme color map for resolving scheme colors to absolute values
        theme_color_map = None
        try:
            from lxml import etree
            for part in prs.part.package.iter_parts():
                if '/ppt/theme/' in str(part.partname):
                    theme = etree.fromstring(part.blob)
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    clr = theme.find('.//a:clrScheme', ns)
                    if clr is not None:
                        theme_color_map = {}
                        for child in clr:
                            tag = child.tag.split('}')[1] if '}' in child.tag else child.tag
                            for c in child:
                                val = c.get('val') or c.get('lastClr')
                                if val:
                                    theme_color_map[tag] = val
        except Exception:
            pass

        # Extract all images once to get filename mapping before metadata
        from ppt2md.parser.image import extract_images_from_slide as extract_imgs
        all_images = extract_imgs(slide, str(images_dir), slide_num, seen_rids)
        img_filename_map = {}
        for img_info in all_images:
            img_filename_map[img_info["shape_index"]] = img_info["filename"]

        shape_metadata = []
        for shape in sorted_shapes:
            meta = extract_shape_metadata(shape, theme_color_map)
            # Inject image filename into metadata for correct roundtrip
            if hasattr(shape, "image") and shape.image is not None:
                idx = list(slide.shapes).index(shape)
                if idx in img_filename_map and "image" in meta:
                    meta["image"]["filename"] = img_filename_map[idx]
            shape_metadata.append(meta)

            shape_md = _process_shape(shape, slide, slide_num, images_dir, media_dir,
                                       seen_rids, formula_as_image, keep_original_format,
                                       verbose, errors)
            if shape_md:
                md_parts.append("\n{}".format(shape_md))

        # Extract formula shapes from AlternateContent (python-pptx skips these)
        from ppt2md.converter.metadata import extract_alternate_content_shapes
        formula_shapes = extract_alternate_content_shapes(slide)
        shape_metadata.extend(formula_shapes)

        # Embed slide metadata as JSON in HTML comment
        import json
        slide_meta = {
            "slide_num": slide_num,
            "shapes": shape_metadata,
        }

        # Store theme color scheme for roundtrip fidelity
        if slide_num == 1:
            try:
                for part in prs.part.package.iter_parts():
                    if '/ppt/theme/' in str(part.partname):
                        theme_xml = part.blob.decode('utf-8')
                        slide_meta["_theme_xml"] = theme_xml
                        break
            except Exception:
                pass
        meta_json = json.dumps(slide_meta, ensure_ascii=False, indent=2)
        md_parts.append("\n<!-- PPTX_META_START\n{}\nPPTX_META_END -->".format(meta_json))

        bg_info = extract_background_image(slide, str(images_dir), slide_num)
        if bg_info:
            md_parts.append("\n{}".format(format_background_markdown(bg_info)))

        media_files = extract_media_from_slide(slide, str(media_dir), slide_num)
        for mf in media_files:
            md_parts.append("\n{}".format(format_media_markdown(mf)))

        ole_objects = detect_ole_objects(slide)
        for ole in ole_objects:
            md_parts.append("\n{}".format(format_ole_markdown(ole)))

        if include_notes:
            notes = extract_notes_from_slide(slide)
            notes_md = format_notes_markdown(notes)
            if notes_md:
                md_parts.append("\n\n{}".format(notes_md))

    output_file = output_dir / "{}.md".format(input_path.stem)
    markdown_content = "\n".join(md_parts) + "\n"

    os.makedirs(output_dir, exist_ok=True)
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(markdown_content)

    result["success"] = True
    result["output_file"] = str(output_file)
    result["errors"] = errors
    return result


def _process_shape(shape, slide, slide_num, images_dir, media_dir,
                    seen_rids, formula_as_image, keep_original_format,
                    verbose, errors):
    """Process a single shape and return Markdown string."""
    parts = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        texts = extract_group_texts(shape)
        return "\n".join(texts) if texts else ""

    if shape.has_text_frame:
        if should_skip_shape(shape):
            return ""

        role = get_placeholder_type(shape)
        if role == "skip":
            return ""

        prefix = placeholder_to_markdown_prefix(role)

        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            level = para.level if para.level else 0
            formatted = format_paragraph_runs(para)
            if not formatted.strip():
                formatted = text

            parts.append("{}{}".format(prefix, formatted))

    if hasattr(shape, "has_table") and shape.has_table:
        tables = extract_tables_from_slide(slide)
        for t in tables:
            if t["shape_index"] == list(slide.shapes).index(shape):
                parts.append(t["markdown"])
                break

    if hasattr(shape, "has_chart") and shape.has_chart:
        charts = extract_charts_from_slide(slide)
        for c in charts:
            if c["shape_index"] == list(slide.shapes).index(shape):
                parts.append(c["markdown"])
                break

        if formula_as_image:
            try:
                chart_images = export_charts_from_slide(slide, str(images_dir), slide_num)
                for ci in chart_images:
                    if ci.get("shape_index") == list(slide.shapes).index(shape):
                        parts.append("![chart](images/{})".format(ci["filename"]))
                        break
            except Exception as e:
                if verbose:
                    handle_image_error(e, list(slide.shapes).index(shape), verbose)

    if hasattr(shape, "image") and shape.image is not None:
        try:
            images = extract_images_from_slide(slide, str(images_dir), slide_num, seen_rids)
            for img in images:
                if img["shape_index"] == list(slide.shapes).index(shape):
                    if keep_original_format:
                        parts.append(format_position_comment(shape))
                    else:
                        parts.append(format_image_markdown(img, shape))
                    break
        except Exception as e:
            handle_image_error(e, list(slide.shapes).index(shape), verbose)

    try:
        omml_elements = find_omml_elements(shape.element)
        for omml in omml_elements:
            try:
                latex = convert_omml_to_latex(omml)
                if latex:
                    parts.append("${}$".format(latex))
            except Exception as e:
                fallback = handle_formula_error(e, list(slide.shapes).index(shape), verbose)
                parts.append(fallback)
    except Exception:
        pass

    return "\n".join(parts)


def main(argv=None):
    parser = argparse.ArgumentParser(
        prog="ppt2md",
        description="Convert PowerPoint (.pptx) files to Markdown format.",
    )
    parser.add_argument(
        "input",
        help="Path to the input .pptx file or directory containing .pptx files",
    )
    parser.add_argument(
        "-o", "--output",
        help="Output directory (default: same as input file)",
    )
    parser.add_argument(
        "--include-notes",
        action="store_true",
        help="Include speaker notes in output",
    )
    parser.add_argument(
        "--include-hidden",
        action="store_true",
        help="Include hidden slides",
    )
    parser.add_argument(
        "--skip-empty",
        action="store_true",
        help="Skip empty slides",
    )
    parser.add_argument(
        "--formula-as-image",
        action="store_true",
        help="Export formulas as images instead of LaTeX",
    )
    parser.add_argument(
        "--keep-original-format",
        action="store_true",
        help="Keep original formatting info in comments",
    )
    parser.add_argument(
        "--image-dpi",
        type=int,
        default=96,
        help="DPI for image export (default: 96)",
    )
    parser.add_argument(
        "--no-frontmatter",
        action="store_true",
        help="Skip generating YAML frontmatter",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Enable verbose output",
    )
    parser.add_argument(
        "--debug",
        action="store_true",
        help="Enable debug output",
    )

    args = parser.parse_args(argv)
    input_path = Path(args.input)

    if not input_path.exists():
        print("Error: {} does not exist".format(input_path), file=sys.stderr)
        return 1

    pptx_files = find_pptx_files(input_path)
    if not pptx_files:
        print("Error: No .pptx files found in {}".format(input_path), file=sys.stderr)
        return 1

    if len(pptx_files) == 1 and input_path.suffix.lower() == ".pptx":
        result = convert_pptx_to_markdown(
            pptx_files[0],
            output_dir=args.output,
            include_notes=args.include_notes,
            include_hidden=args.include_hidden,
            skip_empty=args.skip_empty,
            formula_as_image=args.formula_as_image,
            keep_original_format=args.keep_original_format,
            image_dpi=args.image_dpi,
            no_frontmatter=args.no_frontmatter,
            verbose=args.verbose,
            debug=args.debug,
        )
        if result["success"]:
            print("ppt2md: {} -> {}".format(input_path.name, result["output_file"]))
        else:
            print("ppt2md: Failed to convert {}".format(input_path.name), file=sys.stderr)
        return get_exit_code(result["errors"])

    results = []
    for pptx_file in pptx_files:
        out_dir = args.output if args.output else None
        result = convert_pptx_to_markdown(
            pptx_file,
            output_dir=out_dir,
            include_notes=args.include_notes,
            include_hidden=args.include_hidden,
            skip_empty=args.skip_empty,
            formula_as_image=args.formula_as_image,
            keep_original_format=args.keep_original_format,
            image_dpi=args.image_dpi,
            no_frontmatter=args.no_frontmatter,
            verbose=args.verbose,
            debug=args.debug,
        )
        results.append(result)
        if args.verbose:
            status = "OK" if result["success"] else "FAILED"
            print("  [{}] {} ({} slides)".format(status, pptx_file.name, result["slide_count"]))

    print(batch_convert_summary(results))
    all_errors = []
    for r in results:
        all_errors.extend(r["errors"])
    return get_exit_code(all_errors)


def md_to_pptx_cli(md_path, output_path=None):
    """CLI wrapper for reverse conversion."""
    from ppt2md.converter.reverse import md_to_pptx
    result = md_to_pptx(md_path, output_path=output_path)
    print("ppt2md: {} -> {}".format(md_path, result))
    return 0


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1 and sys.argv[1] == "reverse":
        # ppt2md reverse input.md -o output.pptx
        sys.argv.pop(1)
        parser = argparse.ArgumentParser(prog="ppt2md reverse")
        parser.add_argument("input", help="Path to .md file")
        parser.add_argument("-o", "--output", help="Output .pptx path")
        args = parser.parse_args()
        sys.exit(md_to_pptx_cli(args.input, args.output))
    else:
        sys.exit(main())
