"""Presentation-level PPTX parsing."""

from pptx import Presentation


def open_presentation(file_path):
    """Open a PPTX file and return the Presentation object.

    Args:
        file_path: Path to the .pptx file.

    Returns:
        python-pptx Presentation object.
    """
    return Presentation(file_path)


def get_slide_count(prs):
    """Return the number of slides in the presentation.

    Args:
        prs: Presentation object.

    Returns:
        int: Number of slides.
    """
    return len(prs.slides)


def list_slides(prs):
    """Return basic info about each slide.

    Args:
        prs: Presentation object.

    Returns:
        list of dict with slide index and shape count.
    """
    result = []
    for i, slide in enumerate(prs.slides):
        result.append({
            "index": i,
            "shape_count": len(slide.shapes),
        })
    return result
