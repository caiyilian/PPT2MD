"""SmartArt and group shape handling."""

from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_group_shapes(group_shape):
    """Recursively extract shapes from a group shape.

    Args:
        group_shape: python-pptx GroupShape object.

    Returns:
        list of dict with shape info from the group.
    """
    results = []
    for shape in group_shape.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            results.extend(extract_group_shapes(shape))
        else:
            info = _shape_to_info(shape)
            if info:
                results.append(info)
    return results


def extract_group_texts(group_shape):
    """Extract all text strings from a group shape recursively.

    Args:
        group_shape: python-pptx GroupShape object.

    Returns:
        list of text strings found in the group.
    """
    texts = []
    for shape in group_shape.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            texts.extend(extract_group_texts(shape))
        elif shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
    return texts


def extract_smartart_text(shape):
    """Extract text from a SmartArt shape via XML parsing.

    Args:
        shape: python-pptx Shape that may contain SmartArt.

    Returns:
        list of text strings found in the SmartArt.
    """
    if not hasattr(shape, "element"):
        return []

    texts = []
    ns_map = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }

    for elem in shape.element.iter():
        tag = elem.tag
        if "txBody" in tag or "Body" in tag:
            for para in elem:
                if "p" in para.tag.lower():
                    text_parts = []
                    for run in para:
                        if "r" in run.tag.lower() or "t" in run.tag.lower():
                            if run.text:
                                text_parts.append(run.text)
                        for t_elem in run.iter():
                            if t_elem.text and "t" in t_elem.tag.lower():
                                text_parts.append(t_elem.text)
                    text = "".join(text_parts).strip()
                    if text:
                        texts.append(text)

    return texts


def _shape_to_info(shape):
    """Convert a shape to a dict with relevant info."""
    info = {
        "shape_type": str(shape.shape_type),
        "name": shape.name,
    }

    if shape.has_text_frame:
        texts = []
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                texts.append(text)
        if texts:
            info["texts"] = texts

    if hasattr(shape, "image") and shape.image is not None:
        info["has_image"] = True

    if info.get("texts") or info.get("has_image"):
        return info
    return None
