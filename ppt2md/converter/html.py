"""Render PPTX slides as self-contained HTML."""

import argparse
import base64
import copy
import mimetypes
from html import escape
from pathlib import Path
from tempfile import TemporaryDirectory

from lxml import etree
from pptx import Presentation

from ppt2md.converter.metadata import (
    A_NS,
    P_NS,
    extract_alternate_content_shapes,
    extract_shape_metadata,
)
from ppt2md.parser.image import extract_images_from_slide


EMU_PER_INCH = 914400.0
CSS_DPI = 96.0
EMU_PER_PX = EMU_PER_INCH / CSS_DPI
EMU_PER_PT = 12700.0

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

NS = {
    "a": A_NS,
    "p": P_NS,
    "r": R_NS,
}

THEME_ALIASES = {
    "tx1": "dk1",
    "tx2": "dk2",
    "bg1": "lt1",
    "bg2": "lt2",
}

THEME_ENUM_ALIASES = {
    "DARK_1": "dk1",
    "LIGHT_1": "lt1",
    "DARK_2": "dk2",
    "LIGHT_2": "lt2",
    "TEXT_1": "tx1",
    "TEXT_2": "tx2",
    "BACKGROUND_1": "bg1",
    "BACKGROUND_2": "bg2",
    "ACCENT_1": "accent1",
    "ACCENT_2": "accent2",
    "ACCENT_3": "accent3",
    "ACCENT_4": "accent4",
    "ACCENT_5": "accent5",
    "ACCENT_6": "accent6",
    "HYPERLINK": "hlink",
    "FOLLOWED_HYPERLINK": "folHlink",
}

DASH_ARRAYS = {
    "dash": "8 5",
    "dashDot": "8 4 2 4",
    "dot": "2 4",
    "lgDash": "12 5",
    "lgDashDot": "12 4 2 4",
    "sysDash": "6 3",
    "sysDot": "2 3",
}


class HtmlRenderContext:
    """Rendering state shared across slides and shapes."""

    def __init__(self, theme_color_map, asset_data):
        self.theme_color_map = theme_color_map or {}
        self.asset_data = asset_data or {}
        self._next_id = 0

    def next_id(self, prefix):
        self._next_id += 1
        return "{}-{}".format(prefix, self._next_id)


def convert_pptx_to_html(input_path, output_path=None):
    """Convert a PPTX file to one browser-openable HTML file.

    Args:
        input_path: Source .pptx path.
        output_path: Output .html path. Defaults to source stem + .html.

    Returns:
        Path to the generated HTML file.
    """
    input_path = Path(input_path)
    if output_path is None:
        output_path = input_path.with_suffix(".html")
    output_path = Path(output_path)

    prs = Presentation(str(input_path))

    with TemporaryDirectory() as temp_dir:
        images_dir = Path(temp_dir) / "images"
        slides = _extract_slide_metadata(prs, images_dir)
        asset_data = _load_asset_data(images_dir)
        ctx = HtmlRenderContext(_extract_theme_color_map(prs), asset_data)
        html = _render_document(prs, slides, ctx, input_path.name)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(html, encoding="utf-8")
    return output_path


def _extract_slide_metadata(prs, images_dir):
    """Build metadata blocks for every slide, including extracted images."""
    slides = []
    seen_rids = {}
    theme_color_map = _extract_theme_color_map(prs)

    for slide_num, slide in enumerate(prs.slides, 1):
        images = extract_images_from_slide(slide, str(images_dir), slide_num, seen_rids)
        image_filename_map = {}
        for img_info in images:
            shape_path = img_info.get("shape_path")
            if shape_path is not None:
                key = tuple(shape_path)
            else:
                key = (img_info["shape_index"],)
            image_filename_map[key] = img_info["filename"]
            image_filename_map[(key, img_info.get("kind", "picture"))] = img_info["filename"]
            if img_info.get("rId"):
                image_filename_map[(key, "rId", img_info["rId"])] = img_info["filename"]

        shape_metadata = []
        for idx, shape in enumerate(slide.shapes):
            shape_metadata.append(
                extract_shape_metadata(shape, theme_color_map, image_filename_map, (idx,))
            )
        shape_metadata.extend(extract_alternate_content_shapes(slide))

        slides.append(
            {
                "slide": slide,
                "slide_num": slide_num,
                "shapes": shape_metadata,
                "background": _extract_slide_background(slide, theme_color_map),
            }
        )

    return slides


def _load_asset_data(images_dir):
    """Return filename -> data URI for extracted image assets."""
    assets = {}
    if not images_dir.exists():
        return assets

    for image_path in images_dir.iterdir():
        if not image_path.is_file():
            continue
        mime = mimetypes.guess_type(str(image_path))[0] or "application/octet-stream"
        payload = base64.b64encode(image_path.read_bytes()).decode("ascii")
        assets[image_path.name] = "data:{};base64,{}".format(mime, payload)

    return assets


def _render_document(prs, slides, ctx, source_name):
    slide_width_px = _emu_to_px(prs.slide_width)
    slide_height_px = _emu_to_px(prs.slide_height)
    rendered_slides = []

    for slide_info in slides:
        rendered_slides.append(
            _render_slide(slide_info, ctx, slide_width_px, slide_height_px)
        )

    return """<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{title}</title>
  <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/katex@0.16.11/dist/katex.min.css">
  <script defer src="https://cdn.jsdelivr.net/npm/katex@0.16.11/dist/katex.min.js"></script>
  <script defer src="https://cdn.jsdelivr.net/npm/katex@0.16.11/dist/contrib/auto-render.min.js"
    onload="renderMathInElement(document.body, {{delimiters: [
      {{left: '$$', right: '$$', display: true}},
      {{left: '$', right: '$', display: false}}
    ]}});"></script>
  <style>
{css}
  </style>
</head>
<body>
  <main class="deck" data-source="{source}">
{slides}
  </main>
</body>
</html>
""".format(
        title=escape(Path(source_name).stem),
        source=escape(source_name),
        css=_base_css(),
        slides="\n".join(rendered_slides),
    )


def _base_css():
    return """    * { box-sizing: border-box; }
    body {
      margin: 0;
      background: #f3f4f6;
      color: #111827;
      font-family: Calibri, Arial, sans-serif;
    }
    .deck {
      display: flex;
      flex-direction: column;
      align-items: center;
      gap: 24px;
      padding: 24px;
    }
    .slide {
      position: relative;
      overflow: hidden;
      flex: 0 0 auto;
      background: #fff;
      box-shadow: 0 1px 8px rgba(15, 23, 42, 0.18);
    }
    .ppt-shape,
    .ppt-line,
    .ppt-image,
    .ppt-table {
      position: absolute;
      transform-origin: center center;
    }
    .ppt-shape-svg,
    .ppt-line-svg {
      position: absolute;
      inset: 0;
      width: 100%;
      height: 100%;
      overflow: visible;
    }
    .ppt-shape-text {
      position: absolute;
      display: flex;
      flex-direction: column;
      min-width: 0;
      min-height: 0;
      white-space: pre-wrap;
      overflow-wrap: anywhere;
      line-height: 1.08;
      pointer-events: none;
    }
    .ppt-text-line {
      margin: 0;
      min-height: 1em;
    }
    .ppt-image {
      object-fit: fill;
      display: block;
    }
    .ppt-table table {
      width: 100%;
      height: 100%;
      border-collapse: collapse;
      table-layout: fixed;
      font-family: Calibri, Arial, sans-serif;
      font-size: 12pt;
      color: #111827;
    }
    .ppt-table td {
      border: 1px solid #6b7280;
      padding: 3px 5px;
      vertical-align: middle;
      overflow-wrap: anywhere;
      white-space: pre-wrap;
    }
    @media print {
      body { background: #fff; }
      .deck { display: block; padding: 0; }
      .slide { box-shadow: none; page-break-after: always; margin: 0; }
    }"""


def _render_slide(slide_info, ctx, slide_width_px, slide_height_px):
    style = _style(
        width=slide_width_px,
        height=slide_height_px,
        extra=slide_info["background"],
    )
    shape_html = []
    for z_index, meta in enumerate(slide_info["shapes"], 1):
        shape_html.extend(_render_shape(meta, ctx, z_index))

    return """    <section class="slide" data-slide="{num}" style="{style}">
{shapes}
    </section>""".format(
        num=slide_info["slide_num"],
        style=style,
        shapes="\n".join(shape_html),
    )


def _render_shape(meta, ctx, z_index):
    type_name, type_value = _parse_shape_type(meta.get("type", ""))

    if type_value == 6 or type_name == "GROUP" or meta.get("group"):
        return _render_group(meta, ctx, z_index)
    if type_value == 13 or meta.get("image"):
        return [_render_image(meta, ctx, z_index)]
    if type_value == 9 or type_name == "LINE":
        return [_render_line(meta, ctx, z_index)]
    if meta.get("table"):
        return [_render_table(meta, ctx, z_index)]

    return [_render_auto_shape(meta, ctx, z_index)]


def _render_group(meta, ctx, z_index):
    group_info = meta.get("group") or {}
    children = group_info.get("children", [])
    if not children:
        return []

    group_x = _get_pos(meta, "x")
    group_y = _get_pos(meta, "y")
    group_w = _get_size(meta, "width")
    group_h = _get_size(meta, "height")
    coord_space = group_info.get("coord_space", {})

    rendered = []
    for child_index, child in enumerate(children, 1):
        child_meta = copy.deepcopy(child)
        x, y, w, h = _group_child_bounds_to_slide(
            child_meta, group_x, group_y, group_w, group_h, coord_space
        )
        child_meta["position"] = {"x": x, "y": y}
        child_meta["size"] = {"width": w, "height": h}
        rendered.extend(_render_shape(child_meta, ctx, z_index * 1000 + child_index))
    return rendered


def _render_image(meta, ctx, z_index):
    image_meta = meta.get("image") or {}
    filename = image_meta.get("filename")
    src = ctx.asset_data.get(filename, "") if filename else ""
    if not src:
        src = _transparent_pixel()

    style = _absolute_style(meta, z_index)
    alt = escape(meta.get("name") or "image")
    return '      <img class="ppt-image" alt="{}" src="{}" style="{}">'.format(
        alt, src, style
    )


def _render_line(meta, ctx, z_index):
    line = meta.get("line") or {}
    x = _emu_to_px(_get_pos(meta, "x"))
    y = _emu_to_px(_get_pos(meta, "y"))
    width = max(_emu_to_px(_get_size(meta, "width")), 1.0)
    height = max(_emu_to_px(_get_size(meta, "height")), 1.0)
    rotation = float(meta.get("rotation") or 0)
    stroke = _stroke_info(line, ctx.theme_color_map)
    marker_defs, marker_attrs = _line_markers(line, stroke["color"], ctx)

    flip_h = str(line.get("flipH", "")).lower() in {"1", "true"}
    flip_v = str(line.get("flipV", "")).lower() in {"1", "true"}
    x1 = width if flip_h else 0.0
    y1 = height if flip_v else 0.0
    x2 = 0.0 if flip_h else width
    y2 = 0.0 if flip_v else height

    dash = ' stroke-dasharray="{}"'.format(stroke["dash"]) if stroke["dash"] else ""
    attrs = (
        'stroke="{color}" stroke-width="{stroke_width:.3f}" '
        'stroke-linecap="round" stroke-linejoin="round" fill="none"{dash}{markers}'
    ).format(
        color=stroke["color"],
        stroke_width=stroke["width"],
        dash=dash,
        markers=marker_attrs,
    )

    geom = line.get("connector_geom", "")
    if geom.startswith("bentConnector"):
        mid_x = (x1 + x2) / 2.0
        points = "{:.3f},{:.3f} {:.3f},{:.3f} {:.3f},{:.3f} {:.3f},{:.3f}".format(
            x1, y1, mid_x, y1, mid_x, y2, x2, y2
        )
        primitive = '<polyline points="{}" {}/>'.format(points, attrs)
    else:
        primitive = '<line x1="{:.3f}" y1="{:.3f}" x2="{:.3f}" y2="{:.3f}" {}/>'.format(
            x1, y1, x2, y2, attrs
        )

    style = _style(
        left=x,
        top=y,
        width=width,
        height=height,
        z_index=z_index,
        extra=_rotation_style(rotation),
    )
    return """      <div class="ppt-line" style="{style}">
        <svg class="ppt-line-svg" viewBox="0 0 {width:.3f} {height:.3f}" preserveAspectRatio="none" aria-hidden="true">
          <defs>{defs}</defs>
          {primitive}
        </svg>
      </div>""".format(
        style=style,
        width=width,
        height=height,
        defs=marker_defs,
        primitive=primitive,
    )


def _render_auto_shape(meta, ctx, z_index):
    width = max(_emu_to_px(_get_size(meta, "width")), 1.0)
    height = max(_emu_to_px(_get_size(meta, "height")), 1.0)
    fill_defs, fill = _svg_fill(meta.get("fill"), ctx)
    stroke = _stroke_info(meta.get("line") or {}, ctx.theme_color_map)
    dash = ' stroke-dasharray="{}"'.format(stroke["dash"]) if stroke["dash"] else ""
    stroke_attrs = 'stroke="{}" stroke-width="{:.3f}"{}'.format(
        stroke["color"], stroke["width"], dash
    )
    shape_markup = _shape_svg_markup(
        meta,
        width,
        height,
        fill,
        stroke_attrs,
        ctx,
    )
    text_markup = _render_text(meta, ctx)
    style = _absolute_style(meta, z_index)

    return """      <div class="ppt-shape" style="{style}">
        <svg class="ppt-shape-svg" viewBox="0 0 {width:.3f} {height:.3f}" preserveAspectRatio="none" aria-hidden="true">
          <defs>{defs}</defs>
          {shape}
        </svg>
{text}
      </div>""".format(
        style=style,
        width=width,
        height=height,
        defs=fill_defs,
        shape=shape_markup,
        text=text_markup,
    )


def _render_table(meta, ctx, z_index):
    rows = _parse_table_rows(meta.get("raw_xml"), ctx.theme_color_map)
    if not rows:
        rows = [[{"text": ""} for _ in range((meta.get("table") or {}).get("cols", 1))]]

    html_rows = []
    for row in rows:
        cells = []
        for cell in row:
            attrs = []
            if cell.get("colspan", 1) > 1:
                attrs.append('colspan="{}"'.format(cell["colspan"]))
            if cell.get("rowspan", 1) > 1:
                attrs.append('rowspan="{}"'.format(cell["rowspan"]))
            style = cell.get("style", "")
            if style:
                attrs.append('style="{}"'.format(style))
            cells.append("<td {}>{}</td>".format(" ".join(attrs), escape(cell["text"])))
        html_rows.append("<tr>{}</tr>".format("".join(cells)))

    return """      <div class="ppt-table" style="{style}">
        <table>{rows}</table>
      </div>""".format(
        style=_absolute_style(meta, z_index),
        rows="".join(html_rows),
    )


def _shape_svg_markup(meta, width, height, fill, stroke_attrs, ctx):
    auto_shape_type = (meta.get("auto_shape_type") or "").upper()

    if "OVAL" in auto_shape_type or "ELLIPSE" in auto_shape_type:
        return '<ellipse cx="{cx:.3f}" cy="{cy:.3f}" rx="{rx:.3f}" ry="{ry:.3f}" fill="{fill}" {stroke}/>'.format(
            cx=width / 2.0,
            cy=height / 2.0,
            rx=width / 2.0,
            ry=height / 2.0,
            fill=fill,
            stroke=stroke_attrs,
        )

    if "ROUNDED_RECTANGLE" in auto_shape_type or "ROUNDRECT" in auto_shape_type:
        radius = min(width, height) * 0.16
        return '<rect x="0" y="0" width="{w:.3f}" height="{h:.3f}" rx="{r:.3f}" ry="{r:.3f}" fill="{fill}" {stroke}/>'.format(
            w=width,
            h=height,
            r=radius,
            fill=fill,
            stroke=stroke_attrs,
        )

    if "TRAPEZOID" in auto_shape_type:
        d = "M {:.3f} 0 L {:.3f} 0 L {:.3f} {:.3f} L 0 {:.3f} Z".format(
            width * 0.22, width * 0.78, width, height, height
        )
        return '<path d="{}" fill="{}" {}/>'.format(d, fill, stroke_attrs)

    if "DIAMOND" in auto_shape_type:
        points = "{:.3f},0 {:.3f},{:.3f} {:.3f},{:.3f} 0,{:.3f}".format(
            width / 2.0, width, height / 2.0, width / 2.0, height, height / 2.0
        )
        return '<polygon points="{}" fill="{}" {}/>'.format(points, fill, stroke_attrs)

    if "FLOWCHART_OR" in auto_shape_type:
        return (
            '<ellipse cx="{cx:.3f}" cy="{cy:.3f}" rx="{rx:.3f}" ry="{ry:.3f}" fill="{fill}" {stroke}/>'
            '<line x1="{cx:.3f}" y1="{y1:.3f}" x2="{cx:.3f}" y2="{y2:.3f}" {stroke}/>'
            '<line x1="{x1:.3f}" y1="{cy:.3f}" x2="{x2:.3f}" y2="{cy:.3f}" {stroke}/>'
        ).format(
            cx=width / 2.0,
            cy=height / 2.0,
            rx=width / 2.0,
            ry=height / 2.0,
            x1=width * 0.25,
            x2=width * 0.75,
            y1=height * 0.25,
            y2=height * 0.75,
            fill=fill,
            stroke=stroke_attrs,
        )

    if "FLOWCHART_SUMMING_JUNCTION" in auto_shape_type:
        return (
            '<ellipse cx="{cx:.3f}" cy="{cy:.3f}" rx="{rx:.3f}" ry="{ry:.3f}" fill="{fill}" {stroke}/>'
            '<line x1="{x1:.3f}" y1="{y1:.3f}" x2="{x2:.3f}" y2="{y2:.3f}" {stroke}/>'
            '<line x1="{x2:.3f}" y1="{y1:.3f}" x2="{x1:.3f}" y2="{y2:.3f}" {stroke}/>'
        ).format(
            cx=width / 2.0,
            cy=height / 2.0,
            rx=width / 2.0,
            ry=height / 2.0,
            x1=width * 0.28,
            x2=width * 0.72,
            y1=height * 0.28,
            y2=height * 0.72,
            fill=fill,
            stroke=stroke_attrs,
        )

    if "CUBE" in auto_shape_type:
        inset_x = width * 0.18
        inset_y = height * 0.16
        front = '<rect x="0" y="{iy:.3f}" width="{fw:.3f}" height="{fh:.3f}" fill="{fill}" {stroke}/>'.format(
            iy=inset_y,
            fw=width - inset_x,
            fh=height - inset_y,
            fill=fill,
            stroke=stroke_attrs,
        )
        top = '<polygon points="0,{iy:.3f} {ix:.3f},0 {w:.3f},0 {fw:.3f},{iy:.3f}" fill="{fill}" {stroke}/>'.format(
            iy=inset_y,
            ix=inset_x,
            w=width,
            fw=width - inset_x,
            fill=fill,
            stroke=stroke_attrs,
        )
        side = '<polygon points="{fw:.3f},{iy:.3f} {w:.3f},0 {w:.3f},{h:.3f} {fw:.3f},{h:.3f}" fill="{fill}" {stroke}/>'.format(
            fw=width - inset_x,
            iy=inset_y,
            w=width,
            h=height,
            fill=fill,
            stroke=stroke_attrs,
        )
        return top + side + front

    return '<rect x="0" y="0" width="{w:.3f}" height="{h:.3f}" fill="{fill}" {stroke}/>'.format(
        w=width,
        h=height,
        fill=fill,
        stroke=stroke_attrs,
    )


def _render_text(meta, ctx):
    text_meta = meta.get("text") or {}
    paragraphs = text_meta.get("paragraphs") or []
    if not paragraphs:
        return ""

    body_props = meta.get("body_props") or {}
    inset_left = _emu_to_px(body_props.get("lIns", 45720))
    inset_right = _emu_to_px(body_props.get("rIns", 45720))
    inset_top = _emu_to_px(body_props.get("tIns", 22860))
    inset_bottom = _emu_to_px(body_props.get("bIns", 22860))

    type_name, type_value = _parse_shape_type(meta.get("type", ""))
    default_justify = "flex-start" if type_value == 17 or type_name == "TEXT_BOX" else "center"
    anchor = body_props.get("anchor")
    if anchor in {"b", "bottom"}:
        justify = "flex-end"
    elif anchor in {"ctr", "mid", "center"}:
        justify = "center"
    elif anchor in {"t", "top"}:
        justify = "flex-start"
    else:
        justify = default_justify

    extra = [
        "left:{:.3f}px".format(inset_left),
        "right:{:.3f}px".format(inset_right),
        "top:{:.3f}px".format(inset_top),
        "bottom:{:.3f}px".format(inset_bottom),
        "justify-content:{}".format(justify),
    ]
    if body_props.get("vert"):
        extra.extend(["writing-mode:vertical-rl", "text-orientation:mixed"])

    lines = []
    for para in paragraphs:
        runs = para.get("runs") or []
        if not runs:
            continue
        align = _paragraph_align(para.get("alignment"))
        run_html = "".join(_render_run(run, ctx.theme_color_map) for run in runs)
        lines.append(
            '        <div class="ppt-text-line" style="text-align:{}">{}</div>'.format(
                align, run_html
            )
        )

    if not lines:
        return ""

    return """        <div class="ppt-shape-text" style="{style}">
{lines}
        </div>""".format(
        style=";".join(extra),
        lines="\n".join(lines),
    )


def _render_run(run, theme_color_map):
    text = escape(run.get("text", ""))
    styles = []

    font_size = run.get("font_size")
    if font_size:
        styles.append("font-size:{:.3f}pt".format(float(font_size) / EMU_PER_PT))
    if run.get("font_name"):
        styles.append("font-family:'{}', Calibri, Arial, sans-serif".format(
            escape(str(run["font_name"]), quote=True)
        ))
    color = _css_color(run.get("font_color"), theme_color_map)
    if color:
        styles.append("color:{}".format(color))
    if run.get("bold"):
        styles.append("font-weight:700")
    if run.get("italic"):
        styles.append("font-style:italic")

    decorations = []
    if run.get("underline"):
        decorations.append("underline")
    if run.get("strikethrough"):
        decorations.append("line-through")
    if decorations:
        styles.append("text-decoration:{}".format(" ".join(decorations)))
    if run.get("superscript"):
        styles.extend(["vertical-align:super", "font-size:70%"])
    if run.get("subscript"):
        styles.extend(["vertical-align:sub", "font-size:70%"])

    return '<span style="{}">{}</span>'.format(";".join(styles), text)


def _parse_table_rows(raw_xml, theme_color_map):
    if not raw_xml:
        return []
    try:
        root = etree.fromstring(raw_xml.encode("utf-8"))
    except Exception:
        return []

    tbl = root.find(".//a:tbl", namespaces=NS)
    if tbl is None:
        return []

    parsed_rows = []
    for tr in tbl.findall("a:tr", namespaces=NS):
        parsed_cells = []
        for tc in tr.findall("a:tc", namespaces=NS):
            tc_pr = tc.find("a:tcPr", namespaces=NS)
            if tc_pr is not None and (
                tc_pr.get("hMerge") == "1" or tc_pr.get("vMerge") == "1"
            ):
                continue
            text = "\n".join(
                t.text for t in tc.findall(".//a:t", namespaces=NS) if t.text
            )
            cell = {"text": text}
            if tc_pr is not None:
                if tc_pr.get("gridSpan"):
                    cell["colspan"] = int(tc_pr.get("gridSpan"))
                if tc_pr.get("rowSpan"):
                    cell["rowspan"] = int(tc_pr.get("rowSpan"))
                cell_style = _table_cell_style(tc, theme_color_map)
                if cell_style:
                    cell["style"] = cell_style
            parsed_cells.append(cell)
        parsed_rows.append(parsed_cells)

    return parsed_rows


def _table_cell_style(tc, theme_color_map):
    styles = []
    fill = tc.find(".//a:tcPr/a:solidFill", namespaces=NS)
    color = _color_from_color_parent(fill, theme_color_map)
    if color:
        styles.append("background:{}".format(color))

    p_pr = tc.find(".//a:pPr", namespaces=NS)
    if p_pr is not None and p_pr.get("algn"):
        styles.append("text-align:{}".format(_xml_align(p_pr.get("algn"))))

    return ";".join(styles)


def _extract_theme_color_map(prs):
    color_map = {}
    try:
        for part in prs.part.package.iter_parts():
            if "/ppt/theme/" not in str(part.partname):
                continue
            theme = etree.fromstring(part.blob)
            color_scheme = theme.find(".//a:clrScheme", namespaces=NS)
            if color_scheme is None:
                continue
            for child in color_scheme:
                name = _local_name(child)
                value = None
                for color_child in child:
                    value = color_child.get("val") or color_child.get("lastClr")
                    if value:
                        break
                if value:
                    color_map[name] = value
    except Exception:
        pass

    for alias, target in THEME_ALIASES.items():
        if target in color_map and alias not in color_map:
            color_map[alias] = color_map[target]
    return color_map


def _extract_slide_background(slide, theme_color_map):
    bg_pr = slide._element.find(".//p:cSld/p:bg/p:bgPr", namespaces=NS)
    if bg_pr is None:
        return "background:#fff"

    solid = bg_pr.find("a:solidFill", namespaces=NS)
    if solid is not None:
        color = _color_from_color_parent(solid, theme_color_map)
        if color:
            return "background:{}".format(color)

    grad = bg_pr.find("a:gradFill", namespaces=NS)
    if grad is not None:
        return "background:{}".format(_css_gradient(etree.tostring(grad, encoding="unicode"), theme_color_map))

    return "background:#fff"


def _svg_fill(fill_meta, ctx):
    if not fill_meta:
        return "", "none"

    fill_type = fill_meta.get("type")
    if fill_type == "none":
        return "", "none"

    if fill_type == "solid":
        color = _css_color(fill_meta.get("color"), ctx.theme_color_map, "#ffffff")
        alpha = _alpha_from_xml(fill_meta.get("xml"))
        return "", _with_alpha(color, alpha)

    if fill_type == "scheme":
        color = _css_color(
            fill_meta.get("_resolved") or fill_meta.get("color"),
            ctx.theme_color_map,
            "#ffffff",
        )
        alpha = _alpha_from_xml(fill_meta.get("xml"))
        return "", _with_alpha(color, alpha)

    if fill_type == "gradient":
        gradient_id = ctx.next_id("fill-gradient")
        stops = _gradient_stops(fill_meta.get("xml"), ctx.theme_color_map)
        stop_markup = "".join(
            '<stop offset="{offset:.3f}%" stop-color="{color}" stop-opacity="{opacity:.5f}"/>'.format(
                **stop
            )
            for stop in stops
        )
        defs = '<linearGradient id="{id}" x1="0%" y1="0%" x2="100%" y2="0%">{stops}</linearGradient>'.format(
            id=gradient_id,
            stops=stop_markup,
        )
        return defs, "url(#{})".format(gradient_id)

    if fill_type == "blip" and fill_meta.get("filename"):
        data_uri = ctx.asset_data.get(fill_meta["filename"])
        if data_uri:
            pattern_id = ctx.next_id("fill-image")
            defs = (
                '<pattern id="{id}" patternUnits="objectBoundingBox" width="1" height="1">'
                '<image href="{href}" x="0" y="0" width="1" height="1" preserveAspectRatio="none"/>'
                "</pattern>"
            ).format(id=pattern_id, href=data_uri)
            return defs, "url(#{})".format(pattern_id)

    return "", "none"


def _stroke_info(line_meta, theme_color_map):
    raw_color = line_meta.get("color")
    if raw_color is None and line_meta.get("xml"):
        raw_color = _first_color_from_xml(line_meta["xml"], theme_color_map)
    color = _css_color(raw_color, theme_color_map)
    if not color:
        color = "none"

    width_emu = line_meta.get("width", 0) or 0
    width = _emu_to_px(width_emu) if color != "none" else 0.0
    if color != "none" and width <= 0:
        width = 1.0

    return {
        "color": color,
        "width": width,
        "dash": DASH_ARRAYS.get(line_meta.get("dash"), ""),
    }


def _line_markers(line_meta, color, ctx):
    if color == "none":
        return "", ""

    defs = []
    attrs = []
    if line_meta.get("tailEnd"):
        marker_id = ctx.next_id("tail-arrow")
        defs.append(_marker_def(marker_id, color))
        attrs.append(' marker-start="url(#{})"'.format(marker_id))
    if line_meta.get("headEnd"):
        marker_id = ctx.next_id("head-arrow")
        defs.append(_marker_def(marker_id, color))
        attrs.append(' marker-end="url(#{})"'.format(marker_id))

    return "".join(defs), "".join(attrs)


def _marker_def(marker_id, color):
    return (
        '<marker id="{id}" markerWidth="8" markerHeight="8" refX="7" refY="4" '
        'orient="auto-start-reverse" markerUnits="strokeWidth">'
        '<path d="M 0 0 L 8 4 L 0 8 z" fill="{color}"/></marker>'
    ).format(id=marker_id, color=color)


def _css_gradient(raw_xml, theme_color_map):
    stops = _gradient_stops(raw_xml, theme_color_map)
    if not stops:
        return "#fff"
    css_stops = []
    for stop in stops:
        color = _with_alpha(stop["color"], stop["opacity"])
        css_stops.append("{} {:.3f}%".format(color, stop["offset"]))
    return "linear-gradient(90deg, {})".format(", ".join(css_stops))


def _gradient_stops(raw_xml, theme_color_map):
    if not raw_xml:
        return []
    try:
        root = etree.fromstring(raw_xml.encode("utf-8"))
    except Exception:
        return []

    stops = []
    for gs in root.findall(".//a:gs", namespaces=NS):
        offset = int(gs.get("pos", "0")) / 1000.0
        color = _color_from_color_parent(gs, theme_color_map) or "#ffffff"
        opacity = _alpha_from_node(gs)
        stops.append({"offset": offset, "color": color, "opacity": opacity})
    return stops


def _color_from_color_parent(parent, theme_color_map):
    if parent is None:
        return None
    srgb = parent.find(".//a:srgbClr", namespaces=NS)
    if srgb is not None and srgb.get("val"):
        return _hex_to_css(srgb.get("val"))
    scheme = parent.find(".//a:schemeClr", namespaces=NS)
    if scheme is not None and scheme.get("val"):
        return _css_color(scheme.get("val"), theme_color_map)
    return None


def _first_color_from_xml(raw_xml, theme_color_map):
    try:
        root = etree.fromstring(raw_xml.encode("utf-8"))
    except Exception:
        return None
    color = _color_from_color_parent(root, theme_color_map)
    if color and color.startswith("#"):
        return color[1:]
    return color


def _css_color(value, theme_color_map, default=None):
    if not value:
        return default
    text = str(value)
    if text.startswith("#") and len(text) in {4, 7}:
        return text
    if text.startswith("theme:"):
        text = text[len("theme:") :]
    token = text.split()[0].strip()
    if _is_hex_color(token):
        return _hex_to_css(token)

    scheme = THEME_ENUM_ALIASES.get(token.upper(), token)
    scheme = THEME_ALIASES.get(scheme, scheme)
    if scheme in theme_color_map:
        return _hex_to_css(theme_color_map[scheme])
    return default


def _hex_to_css(value):
    if not value:
        return None
    value = str(value).strip().lstrip("#")
    if len(value) == 6:
        return "#{}".format(value)
    return None


def _with_alpha(color, alpha):
    if alpha is None or alpha >= 0.999 or not color or not color.startswith("#") or len(color) != 7:
        return color
    r = int(color[1:3], 16)
    g = int(color[3:5], 16)
    b = int(color[5:7], 16)
    return "rgba({}, {}, {}, {:.5f})".format(r, g, b, alpha)


def _alpha_from_xml(raw_xml):
    if not raw_xml:
        return None
    try:
        root = etree.fromstring(raw_xml.encode("utf-8"))
    except Exception:
        return None
    return _alpha_from_node(root)


def _alpha_from_node(node):
    alpha = node.find(".//a:alpha", namespaces=NS)
    if alpha is None or alpha.get("val") is None:
        return 1.0
    try:
        return max(0.0, min(1.0, int(alpha.get("val")) / 100000.0))
    except ValueError:
        return 1.0


def _absolute_style(meta, z_index):
    return _style(
        left=_emu_to_px(_get_pos(meta, "x")),
        top=_emu_to_px(_get_pos(meta, "y")),
        width=max(_emu_to_px(_get_size(meta, "width")), 1.0),
        height=max(_emu_to_px(_get_size(meta, "height")), 1.0),
        z_index=z_index,
        extra=_rotation_style(float(meta.get("rotation") or 0)),
    )


def _style(left=None, top=None, width=None, height=None, z_index=None, extra=None):
    parts = []
    if left is not None:
        parts.append("left:{:.3f}px".format(left))
    if top is not None:
        parts.append("top:{:.3f}px".format(top))
    if width is not None:
        parts.append("width:{:.3f}px".format(width))
    if height is not None:
        parts.append("height:{:.3f}px".format(height))
    if z_index is not None:
        parts.append("z-index:{}".format(z_index))
    if extra:
        parts.append(extra.rstrip(";"))
    return ";".join(parts)


def _rotation_style(rotation):
    if not rotation:
        return ""
    return "transform:rotate({:.3f}deg)".format(rotation)


def _paragraph_align(value):
    if not value:
        return "left"
    if "CENTER" in value:
        return "center"
    if "RIGHT" in value:
        return "right"
    if "JUSTIFY" in value:
        return "justify"
    return "left"


def _xml_align(value):
    return {
        "ctr": "center",
        "r": "right",
        "just": "justify",
        "dist": "justify",
    }.get(value, "left")


def _group_child_bounds_to_slide(child_meta, group_x, group_y, group_w, group_h, coord_space):
    child_x = _get_pos(child_meta, "x")
    child_y = _get_pos(child_meta, "y")
    child_w = _get_size(child_meta, "width")
    child_h = _get_size(child_meta, "height")

    ch_off_x = coord_space.get("chOffX", group_x)
    ch_off_y = coord_space.get("chOffY", group_y)
    ch_ext_w = coord_space.get("chExtCX", group_w) or group_w or 1
    ch_ext_h = coord_space.get("chExtCY", group_h) or group_h or 1
    scale_x = float(group_w) / float(ch_ext_w)
    scale_y = float(group_h) / float(ch_ext_h)

    return (
        int(group_x + (child_x - ch_off_x) * scale_x),
        int(group_y + (child_y - ch_off_y) * scale_y),
        int(child_w * scale_x),
        int(child_h * scale_y),
    )


def _get_pos(meta, key):
    return int((meta.get("position") or {}).get(key, 0) or 0)


def _get_size(meta, key):
    return int((meta.get("size") or {}).get(key, 0) or 0)


def _emu_to_px(value):
    return float(value or 0) / EMU_PER_PX


def _parse_shape_type(type_text):
    if not type_text:
        return "", 0
    name = type_text.split("(", 1)[0].strip()
    try:
        value = int(type_text.rsplit("(", 1)[1].split(")", 1)[0])
    except (IndexError, ValueError):
        value = 0
    return name, value


def _is_hex_color(value):
    if not value:
        return False
    value = str(value).strip().lstrip("#")
    if len(value) != 6:
        return False
    try:
        int(value, 16)
        return True
    except ValueError:
        return False


def _local_name(node):
    if "}" in node.tag:
        return node.tag.rsplit("}", 1)[1]
    return node.tag


def _transparent_pixel():
    return (
        "data:image/gif;base64,"
        "R0lGODlhAQABAIAAAAAAAP///ywAAAAAAQABAAACAUwAOw=="
    )


def screenshot_html(html_path, output_png=None):
    """Render a local HTML file to a PNG screenshot using Playwright.

    Args:
        html_path: Path to the HTML file.
        output_png: Output PNG path (default: html_path with .png suffix).

    Returns:
        Path to the generated PNG file.
    """
    html_path = Path(html_path)
    if output_png is None:
        output_png = html_path.with_suffix(".png")
    output_png = Path(output_png)

    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        print("  [WARN] playwright not installed. Run: pip install playwright && python -m playwright install chromium")
        return None

    file_url = "file:///" + str(html_path.resolve()).replace("\\", "/")
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page()
        page.goto(file_url, wait_until="networkidle")
        page.screenshot(path=str(output_png.resolve()), full_page=True)
        browser.close()

    return output_png


def main(argv=None):
    parser = argparse.ArgumentParser(
        prog="ppt2html.py",
        description="Convert a PPTX file to a self-contained HTML reconstruction.",
    )
    parser.add_argument("input", help="Path to the input .pptx file")
    parser.add_argument("-o", "--output", help="Output HTML file path")
    parser.add_argument("--screenshot", action="store_true",
                        help="Also take a Playwright screenshot of the generated HTML for visual comparison")
    args = parser.parse_args(argv)

    input_path = Path(args.input)
    if not input_path.exists():
        parser.error("{} does not exist".format(input_path))

    output_path = convert_pptx_to_html(input_path, args.output)
    print("ppt2html: {} -> {}".format(input_path, output_path))

    if args.screenshot:
        screenshot_path = output_path.with_suffix(".png")
        screenshot_html(str(output_path), str(screenshot_path))
        print("ppt2html: screenshot -> {}".format(screenshot_path))

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
