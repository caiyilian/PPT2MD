"""Convert Markdown with embedded metadata back to PPTX."""

import json
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from lxml import etree

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def parse_metadata_blocks(md_content, tag="PPTX_META"):
    """Extract JSON metadata blocks from markdown content.

    Args:
        md_content: Markdown content string.
        tag: Tag name to look for (default: PPTX_META).

    Returns list of parsed dicts.
    """
    pattern = r'<!-- {}_START\n(.*?)\n{}_END -->'.format(tag, tag)
    matches = re.findall(pattern, md_content, re.DOTALL)
    results = []
    for match in matches:
        try:
            results.append(json.loads(match))
        except json.JSONDecodeError:
            continue
    return results


def parse_frontmatter(md_content):
    """Extract YAML frontmatter from markdown."""
    match = re.match(r'^---\n(.*?)\n---', md_content, re.DOTALL)
    if not match:
        return {}
    fm = {}
    for line in match.group(1).split('\n'):
        if ':' in line:
            key, _, value = line.partition(':')
            value = value.strip().strip('"')
            fm[key.strip()] = value
    return fm


def _parse_shape_type_str(type_str):
    """Parse shape type string like 'AUTO_SHAPE (1)' or 'LINE (9)'.

    Returns (name, int_value) or None.
    """
    if not type_str:
        return None
    m = re.match(r'(\w+)\s*\((\d+)\)', type_str)
    if m:
        return m.group(1), int(m.group(2))
    return None


def md_to_pptx(md_path, images_dir=None, output_path=None):
    """Convert a markdown file with embedded metadata back to PPTX.

    Args:
        md_path: Path to the .md file.
        images_dir: Path to images directory (default: same dir as md).
        output_path: Output .pptx path (default: same name as md).

    Returns:
        Path to the created PPTX file.
    """
    md_path = Path(md_path)
    if images_dir is None:
        images_dir = md_path.parent / "images"
    else:
        images_dir = Path(images_dir)
    if output_path is None:
        output_path = md_path.with_suffix(".pptx")
    else:
        output_path = Path(output_path)

    with open(md_path, "r", encoding="utf-8") as f:
        md_content = f.read()

    frontmatter = parse_frontmatter(md_content)
    slides_meta = parse_metadata_blocks(md_content, "PPTX_META")

    pres_meta_list = parse_metadata_blocks(md_content, "PPTX_PRESENTATION_META")
    pres_meta = pres_meta_list[0] if pres_meta_list else {}

    prs = Presentation()
    prs.slide_width = Emu(int(pres_meta.get("slide_width", 12192000)))
    prs.slide_height = Emu(int(pres_meta.get("slide_height", 6858000)))

    # Set document properties
    if frontmatter.get("title"):
        prs.core_properties.title = frontmatter["title"]
    if frontmatter.get("author"):
        prs.core_properties.author = frontmatter["author"]

    slide_layout = prs.slide_layouts[6]  # Blank layout

    # Remove objectDefaults from theme (contains style/fontRef=lt1 that
    # causes Office to mis-resolve schemeClr colors, same as per-shape p:style)
    _remove_theme_object_defaults(prs)

    for slide_meta in slides_meta:
        slide = prs.slides.add_slide(slide_layout)

        shapes_meta = slide_meta.get("shapes", [])
        for shape_meta in shapes_meta:
            _add_shape_from_metadata(slide, shape_meta, images_dir)

    prs.save(str(output_path))

    # Post-save: replace theme XML directly in ZIP to avoid python-pptx caching issues
    for slide_meta in slides_meta:
        theme_xml_str = slide_meta.get("_theme_xml", "")
        if theme_xml_str:
            try:
                import zipfile, io
                with zipfile.ZipFile(str(output_path), 'r') as zin:
                    data = {item.filename: zin.read(item.filename) for item in zin.infolist()}
                for key in list(data.keys()):
                    if 'ppt/theme/' in key:
                        data[key] = theme_xml_str.encode('utf-8')
                with zipfile.ZipFile(str(output_path), 'w', zipfile.ZIP_DEFLATED) as zout:
                    for name, content in data.items():
                        zout.writestr(name, content)
            except Exception as e:
                import traceback
                traceback.print_exc()
        break

    return output_path


def convert_md_to_pptx(md_path, output_path=None, images_dir=None):
    """Compatibility wrapper for MD to PPTX roundtrip conversion.

    The public issue workflow and compare_roundtrip.py use this function name,
    while the original implementation exposed md_to_pptx().
    """
    return md_to_pptx(md_path, images_dir=images_dir, output_path=output_path)


def _add_shape_from_metadata(slide, meta, images_dir):
    """Add a shape to a slide based on metadata."""
    parsed = _parse_shape_type_str(meta.get("type", ""))
    type_name = parsed[0] if parsed else ""
    type_val = parsed[1] if parsed else 0

    if _add_raw_shape_xml(slide, meta.get("raw_xml"), meta.get("raw_relationships"), images_dir):
        return

    x = meta.get("position", {}).get("x", 0)
    y = meta.get("position", {}).get("y", 0)
    w = meta.get("size", {}).get("width", 100000)
    h = meta.get("size", {}).get("height", 100000)
    rotation = meta.get("rotation", 0)

    # Handle images (MSO_SHAPE_TYPE.PICTURE = 13)
    if type_val == 13 or ("image" in meta and meta["image"]):
        _add_image_shape(slide, meta, images_dir, x, y, w, h)
        return

    # Handle tables
    if "table" in meta and meta["table"]:
        _add_table_shape(slide, meta, x, y, w, h)
        return

    # Handle charts (skip - complex to recreate)
    if "chart" in meta and meta["chart"]:
        return

    # Handle lines/arrows (MSO_SHAPE_TYPE.LINE = 9)
    if type_val == 9 or type_name == "LINE":
        _add_line_shape(slide, meta, x, y, w, h, rotation)
        return

    # Handle groups (MSO_SHAPE_TYPE.GROUP = 6)
    if type_val == 6 or type_name == "GROUP":
        _add_group_shape(slide, meta, images_dir, x, y, w, h)
        return

    # Handle formula shapes (AlternateContent with OMML)
    if meta.get("_is_formula"):
        _add_formula_shape(slide, meta, x, y, w, h, rotation)
        return

    # Handle text boxes (MSO_SHAPE_TYPE.TEXT_BOX = 17)
    if type_val == 17 or type_name == "TEXT_BOX":
        _add_text_box_shape(slide, meta, images_dir, x, y, w, h, rotation)
        return

    # Handle auto shapes and everything else
    _add_auto_shape(slide, meta, images_dir, x, y, w, h, rotation)


def _add_group_shape(slide, meta, images_dir, x, y, w, h):
    """Add a group shape with children."""
    try:
        P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        from pptx.oxml.ns import qn

        group_info = meta.get("group")
        if not group_info or "children" not in group_info:
            return

        # Create group element
        grpSp = etree.SubElement(slide.shapes._spTree, qn('p:grpSp'))

        # Get a unique shape ID from the spTree
        nvGrpSpPr = slide.shapes._spTree.find(qn('p:nvGrpSpPr'))
        next_id = 1
        if nvGrpSpPr is not None:
            existing = nvGrpSpPr.findall('.//' + qn('p:cNvPr'))
            for el in existing:
                try:
                    val = int(el.get('id', '0'))
                    if val >= next_id:
                        next_id = val + 1
                except:
                    pass

        # nvGrpSpPr
        nv = etree.SubElement(grpSp, qn('p:nvGrpSpPr'))
        cnvPr = etree.SubElement(nv, qn('p:cNvPr'))
        cnvPr.set('id', str(next_id))
        next_id += 1
        cnvPr.set('name', meta.get("name", "Group"))
        etree.SubElement(nv, qn('p:cNvGrpSpPr'))
        etree.SubElement(nv, qn('p:nvPr'))

        # grpSpPr with xfrm
        grpSpPr = etree.SubElement(grpSp, qn('p:grpSpPr'))
        xfrm = etree.SubElement(grpSpPr, qn('a:xfrm'))
        off = etree.SubElement(xfrm, qn('a:off'))
        off.set('x', str(x))
        off.set('y', str(y))
        ext = etree.SubElement(xfrm, qn('a:ext'))
        ext.set('cx', str(w))
        ext.set('cy', str(h))

        # Coordinate space for children
        cs = group_info.get("coord_space", {})
        chOff = etree.SubElement(xfrm, qn('a:chOff'))
        chOff.set('x', str(cs.get('chOffX', 0)))
        chOff.set('y', str(cs.get('chOffY', 0)))
        chExt = etree.SubElement(xfrm, qn('a:chExt'))
        chExt.set('cx', str(cs.get('chExtCX', w)))
        chExt.set('cy', str(cs.get('chExtCY', h)))

        # Add each child shape
        for idx, child_meta in enumerate(group_info["children"]):
            _build_group_child(grpSp, child_meta, next_id, slide, images_dir)
            next_id += 1

    except Exception:
        pass


def _add_raw_shape_xml(slide, raw_xml, raw_relationships=None, images_dir=None):
    """Append raw shape XML when it has no external relationships to remap."""
    if not raw_xml:
        return False
    try:
        el = etree.fromstring(raw_xml.encode("utf-8"))
        if not _remap_raw_relationships(slide, el, raw_relationships or {}, images_dir):
            return False
        slide.shapes._spTree.append(el)
        return True
    except Exception:
        return False


def _remap_raw_relationships(slide, el, raw_relationships, images_dir):
    """Recreate and replace rIds referenced by raw shape XML."""
    rel_attrs = [
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link",
    ]
    referenced = []
    for node in el.iter():
        for attr in rel_attrs:
            rid = node.get(attr)
            if rid:
                referenced.append((node, attr, rid))

    if not referenced:
        return True
    if not images_dir:
        return False

    for node, attr, old_rid in referenced:
        rel = raw_relationships.get(old_rid)
        if not rel or rel.get("type") != "image":
            return False
        img_path = Path(images_dir) / rel.get("filename", "")
        if not img_path.exists():
            return False
        _, new_rid = slide.part.get_or_add_image_part(str(img_path))
        node.set(attr, new_rid)

    return True


def _group_child_bounds_to_slide(child_meta, group_x, group_y, group_w, group_h, coord_space):
    """Map a group child bounding box from group coordinates to slide coordinates."""
    child_x = child_meta.get("position", {}).get("x", 0)
    child_y = child_meta.get("position", {}).get("y", 0)
    child_w = child_meta.get("size", {}).get("width", 100000)
    child_h = child_meta.get("size", {}).get("height", 100000)

    ch_off_x = coord_space.get("chOffX", group_x)
    ch_off_y = coord_space.get("chOffY", group_y)
    ch_ext_w = coord_space.get("chExtCX", group_w) or group_w or 1
    ch_ext_h = coord_space.get("chExtCY", group_h) or group_h or 1
    scale_x = float(group_w) / float(ch_ext_w)
    scale_y = float(group_h) / float(ch_ext_h)

    return (
        int(group_x + (child_x - ch_off_x) * scale_x),
        int(group_y + (child_y - ch_off_y) * scale_y),
        int(child_w * scale_x),
        int(child_h * scale_y),
    )


def _add_formula_shape(slide, meta, x, y, w, h, rotation):
    """Add a formula shape wrapped in AlternateContent with OMML XML."""
    from pptx.oxml.ns import qn
    from lxml import etree
    import re

    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    MC_NS = 'http://schemas.openxmlformats.org/markup-compatibility/2006'
    A14_NS = 'http://schemas.microsoft.com/office/drawing/2010/main'
    R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    omml_xml_list = meta.get("_omml_xml", [])
    fallback_xml = meta.get("_fallback_xml")

    if not omml_xml_list:
        # Fallback: create regular text box with LaTeX text
        _add_text_box_shape(slide, meta, None, x, y, w, h, rotation)
        return

    try:
        from pptx.oxml import parse_xml

        # Get a unique shape ID
        nv_container = slide.shapes._spTree.find(qn('p:nvGrpSpPr'))
        sp_id = 1
        if nv_container is not None:
            for el in nv_container.iter():
                if el.tag.endswith('}cNvPr'):
                    try:
                        val = int(el.get('id', '0'))
                        if val >= sp_id:
                            sp_id = val + 1
                    except:
                        pass

        # Build the AlternateContent structure as XML string
        # Use parse_xml to avoid lxml namespace registration conflicts

        # Choice sp XML (without OMML, which we inject separately)
        choice_sp_xml = """<p:sp xmlns:p="%s" xmlns:a="%s" xmlns:r="%s">
  <p:nvSpPr>
    <p:cNvPr id="%d" name="%s"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="none" rtlCol="0"><a:spAutoFit/></a:bodyPr>
    <a:lstStyle/>
    <a:p><a:pPr/>
    </a:p>
  </p:txBody>
</p:sp>""" % (P_NS, A_NS, R_NS, sp_id, meta.get('name', 'Formula'), x, y, w, h)
        sp_id += 1

        choice_sp = etree.fromstring(choice_sp_xml.encode())

        # Add regular text runs from metadata (before OMML elements)
        p_el = choice_sp.find('.//{%s}p' % A_NS)
        text_meta = meta.get('text', {})
        for para in text_meta.get('paragraphs', []):
            for run_meta in para.get('runs', []):
                text = run_meta.get('text', '')
                # Skip formula text (starts with $) - it's handled by OMML
                if text.startswith('$') and text.endswith('$'):
                    continue
                if not text:
                    continue
                # Add as a regular a:r element
                r = etree.SubElement(p_el, qn('a:r'))
                rPr = etree.SubElement(r, qn('a:rPr'))
                if run_meta.get('font_size'):
                    rPr.set('sz', str(int(run_meta['font_size'] * 100 / 12700)))
                if run_meta.get('bold'):
                    rPr.set('b', '1')
                if run_meta.get('superscript'):
                    rPr.set('baseline', '30000')
                if run_meta.get('subscript'):
                    rPr.set('baseline', '-25000')
                t = etree.SubElement(r, qn('a:t'))
                t.text = text

        # Inject OMML XML into the a:p element (after regular text)
        for omml_xml in omml_xml_list:
            omml_el = etree.fromstring(omml_xml.encode())
            if p_el is not None:
                p_el.append(omml_el)

        # Fallback XML
        fb_xml = """<p:sp xmlns:p="%s" xmlns:a="%s">
  <p:nvSpPr>
    <p:cNvPr id="%d" name="%s"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:r><a:rPr/><a:t> </a:t></a:r></a:p>
  </p:txBody>
</p:sp>""" % (P_NS, A_NS, sp_id, meta.get('name', 'Formula') + '_fb', x, y, w, h)
        sp_id += 1

        # Full AlternateContent XML
        mc_xml = """<mc:AlternateContent xmlns:mc="%s" xmlns:a14="%s"
  xmlns:a="%s" xmlns:p="%s" xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math"
  xmlns:r="%s">
  <mc:Choice Requires="a14">
    %s
  </mc:Choice>
  <mc:Fallback>
    %s
  </mc:Fallback>
</mc:AlternateContent>""" % (MC_NS, A14_NS, A_NS, P_NS, R_NS,
                               etree.tostring(choice_sp).decode(),
                               fb_xml)

        ac_element = parse_xml(mc_xml.encode())
        slide.shapes._spTree.append(ac_element)

    except Exception:
        # Fallback to text box
        _add_text_box_shape(slide, meta, None, x, y, w, h, rotation)


def _build_group_child(grpSp, meta, idx, slide=None, images_dir=None):
    """Build a child shape element inside a group."""
    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    from pptx.oxml.ns import qn

    type_str = meta.get("type", "")
    type_val = _parse_type_value(type_str)

    x = meta.get("position", {}).get("x", 0)
    y = meta.get("position", {}).get("y", 0)
    w = meta.get("size", {}).get("width", 100000)
    h = meta.get("size", {}).get("height", 100000)
    rotation = meta.get("rotation", 0)
    name = meta.get("name", "Shape %d" % idx)

    if type_val == 6 or "group" in meta:
        _build_nested_group(grpSp, meta, idx, slide, images_dir)
        return

    if type_val == 13 or meta.get("image"):
        _build_group_picture(grpSp, meta, idx, slide, images_dir)
        return

    if type_val == 9:  # LINE / CONNECTOR
        el = etree.SubElement(grpSp, qn('p:cxnSp'))
        nv = etree.SubElement(el, qn('p:nvCxnSpPr'))
        cnvPr = etree.SubElement(nv, qn('p:cNvPr'))
        cnvPr.set('id', str(idx))
        cnvPr.set('name', name)
        etree.SubElement(nv, qn('p:cNvCxnSpPr'))
        etree.SubElement(nv, qn('p:nvPr'))
    else:
        el = etree.SubElement(grpSp, qn('p:sp'))
        nv = etree.SubElement(el, qn('p:nvSpPr'))
        cnvPr = etree.SubElement(nv, qn('p:cNvPr'))
        cnvPr.set('id', str(idx))
        cnvPr.set('name', name)
        etree.SubElement(nv, qn('p:cNvSpPr'))
        etree.SubElement(nv, qn('p:nvPr'))

    # spPr
    spPr = etree.SubElement(el, qn('p:spPr'))
    xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    if rotation:
        xfrm.set('rot', str(int(rotation * 60000)))
    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(x))
    off.set('y', str(y))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(w))
    ext.set('cy', str(h))

    if type_val != 9:  # Not a connector - add geometry
        auto_type_str = meta.get("auto_shape_type", "")
        prst = 'rect'
        if auto_type_str:
            parsed = re.match(r'.*?\((\d+)\)', auto_type_str)
            if parsed:
                auto_val = int(parsed.group(1))
                try:
                    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                    prst = MSO_AUTO_SHAPE_TYPE(auto_val).xml_value
                except Exception:
                    prst_map = {1: 'rect', 5: 'roundRect', 9: 'ellipse', 33: 'rightArrow', 34: 'leftArrow', 77: 'flowChartSummingJunction', 78: 'flowChartOr'}
                    prst = prst_map.get(auto_val, 'rect')
        geom = etree.SubElement(spPr, qn('a:prstGeom'))
        geom.set('prst', prst)
        etree.SubElement(geom, qn('a:avLst'))
    else:
        # Connector: add line geometry
        geom = etree.SubElement(spPr, qn('a:prstGeom'))
        geom.set('prst', 'line')
        etree.SubElement(geom, qn('a:avLst'))

    # Fill (skip for connectors - connectors use line fill only)
    fill_meta = meta.get("fill")
    if fill_meta and type_val != 9:
        _apply_fill_xml(spPr, fill_meta)

    # Line
    line_meta = meta.get("line")
    if line_meta:
        _apply_line_xml(spPr, line_meta)

    # Text
    text_meta = meta.get("text")
    if text_meta:
        _add_body_props_xml(el, text_meta)
        _apply_text_xml(el, text_meta)

    # Body props
    body_props = meta.get("body_props")
    if body_props:
        _apply_body_props_xml(el, body_props)


def _build_nested_group(parent, meta, idx, slide=None, images_dir=None):
    """Build a nested p:grpSp element."""
    from pptx.oxml.ns import qn

    group_info = meta.get("group") or {}
    x = meta.get("position", {}).get("x", 0)
    y = meta.get("position", {}).get("y", 0)
    w = meta.get("size", {}).get("width", 100000)
    h = meta.get("size", {}).get("height", 100000)

    grp = etree.SubElement(parent, qn('p:grpSp'))
    nv = etree.SubElement(grp, qn('p:nvGrpSpPr'))
    cnvPr = etree.SubElement(nv, qn('p:cNvPr'))
    cnvPr.set('id', str(idx))
    cnvPr.set('name', meta.get("name", "Group %d" % idx))
    etree.SubElement(nv, qn('p:cNvGrpSpPr'))
    etree.SubElement(nv, qn('p:nvPr'))

    grpSpPr = etree.SubElement(grp, qn('p:grpSpPr'))
    xfrm = etree.SubElement(grpSpPr, qn('a:xfrm'))
    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(x))
    off.set('y', str(y))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(w))
    ext.set('cy', str(h))

    cs = group_info.get("coord_space", {})
    chOff = etree.SubElement(xfrm, qn('a:chOff'))
    chOff.set('x', str(cs.get('chOffX', x)))
    chOff.set('y', str(cs.get('chOffY', y)))
    chExt = etree.SubElement(xfrm, qn('a:chExt'))
    chExt.set('cx', str(cs.get('chExtCX', w)))
    chExt.set('cy', str(cs.get('chExtCY', h)))

    next_id = idx * 1000
    for child in group_info.get("children", []):
        _build_group_child(grp, child, next_id, slide, images_dir)
        next_id += 1


def _build_group_picture(parent, meta, idx, slide=None, images_dir=None):
    """Build a p:pic element inside a group."""
    from pptx.oxml.ns import qn

    image_meta = meta.get("image", {})
    filename = image_meta.get("filename")
    if not slide or not images_dir or not filename:
        return

    img_path = Path(images_dir) / filename
    if not img_path.exists():
        return

    _, r_id = slide.part.get_or_add_image_part(str(img_path))

    x = meta.get("position", {}).get("x", 0)
    y = meta.get("position", {}).get("y", 0)
    w = meta.get("size", {}).get("width", 100000)
    h = meta.get("size", {}).get("height", 100000)
    name = meta.get("name", "Picture %d" % idx)

    pic = etree.SubElement(parent, qn('p:pic'))
    nv = etree.SubElement(pic, qn('p:nvPicPr'))
    cnvPr = etree.SubElement(nv, qn('p:cNvPr'))
    cnvPr.set('id', str(idx))
    cnvPr.set('name', name)
    etree.SubElement(nv, qn('p:cNvPicPr'))
    etree.SubElement(nv, qn('p:nvPr'))

    blipFill = etree.SubElement(pic, qn('p:blipFill'))
    blip = etree.SubElement(blipFill, qn('a:blip'))
    blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', r_id)
    stretch = etree.SubElement(blipFill, qn('a:stretch'))
    etree.SubElement(stretch, qn('a:fillRect'))

    spPr = etree.SubElement(pic, qn('p:spPr'))
    xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(x))
    off.set('y', str(y))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(w))
    ext.set('cy', str(h))
    geom = etree.SubElement(spPr, qn('a:prstGeom'))
    geom.set('prst', 'rect')
    etree.SubElement(geom, qn('a:avLst'))


def _apply_fill_xml(spPr, fill_meta):
    """Apply fill to an spPr XML element (used for group children)."""
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    from pptx.oxml.ns import qn

    raw_xml = fill_meta.get("xml")
    if raw_xml:
        _apply_raw_fill_to_spPr(spPr, raw_xml)
        return

    fill_type = fill_meta.get("type")
    if fill_type == "none":
        etree.SubElement(spPr, qn('a:noFill'))
        return

    if fill_type == "scheme":
        # Use resolved absolute color if available
        resolved = fill_meta.get("_resolved")
        if resolved and _is_hex_color(resolved):
            solid = etree.SubElement(spPr, qn('a:solidFill'))
            sc = etree.SubElement(solid, qn('a:srgbClr'))
            sc.set('val', resolved)
            return
        color = fill_meta.get("color", "bg1")
        modifiers = fill_meta.get("modifiers")
        solid = etree.SubElement(spPr, qn('a:solidFill'))
        sc = etree.SubElement(solid, qn('a:schemeClr'))
        sc.set('val', color)
        if modifiers:
            for k, v in modifiers.items():
                mod = etree.SubElement(sc, '{%s}%s' % (A_NS, k))
                mod.set('val', v)
        return

    if fill_type == "rgb":
        color = fill_meta.get("color", "000000")
        solid = etree.SubElement(spPr, qn('a:solidFill'))
        sc = etree.SubElement(solid, qn('a:srgbClr'))
        sc.set('val', color.upper())
        return

    if fill_type in ("blip", "picture"):
        # Skip image fills for group children (complex)
        pass


def _apply_line_xml(spPr, line_meta):
    """Apply line to an spPr XML element (used for group children)."""
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    from pptx.oxml.ns import qn

    raw_xml = line_meta.get("xml")
    if raw_xml:
        _apply_raw_line_to_spPr(spPr, raw_xml)
        return

    width = line_meta.get("width", 0)
    color = line_meta.get("color")

    if color is None and width == 0:
        return

    ln = etree.SubElement(spPr, qn('a:ln'))
    if width > 0:
        ln.set('w', str(width))

    if color is None:
        return

    solid = etree.SubElement(ln, qn('a:solidFill'))
    if _is_hex_color(color):
        sc = etree.SubElement(solid, qn('a:srgbClr'))
        sc.set('val', color.upper())
    else:
        sc = etree.SubElement(solid, qn('a:schemeClr'))
        sc.set('val', color)


def _add_body_props_xml(el, text_meta):
    """Add txBody with bodyPr to a shape element."""
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    from pptx.oxml.ns import qn

    txBody = etree.SubElement(el, qn('p:txBody'))
    bodyPr = etree.SubElement(txBody, qn('a:bodyPr'))
    bodyPr.set('rtlCol', '0')
    bodyPr.set('anchor', 'ctr')
    etree.SubElement(txBody, qn('a:lstStyle'))


def _apply_text_xml(el, text_meta):
    """Apply text to a group child element."""
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    from pptx.oxml.ns import qn

    paragraphs = text_meta.get("paragraphs", [])
    if not paragraphs:
        return

    # Find or create txBody
    txBody = el.find(qn('p:txBody'))
    if txBody is None:
        txBody = etree.SubElement(el, qn('p:txBody'))
        bodyPr = etree.SubElement(txBody, qn('a:bodyPr'))
        bodyPr.set('rtlCol', '0')
        etree.SubElement(txBody, qn('a:lstStyle'))

    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None and bodyPr.get('anchor') is None:
        bodyPr.set('anchor', 'ctr')

    for para_meta in paragraphs:
        p = etree.SubElement(txBody, qn('a:p'))
        level = para_meta.get("level", 0)
        if level:
            pPr = etree.SubElement(p, qn('a:pPr'))
            pPr.set('lvl', str(level))

        for run_meta in para_meta.get("runs", []):
            r = etree.SubElement(p, qn('a:r'))
            text = run_meta.get("text", "")

            # Build rPr
            rPr = etree.SubElement(r, qn('a:rPr'))
            font_size = run_meta.get("font_size")
            if font_size:
                rPr.set('sz', str(int(font_size * 100 / 12700)))
            if run_meta.get("bold"):
                rPr.set('b', '1')
            if run_meta.get("italic"):
                rPr.set('i', '1')
            if run_meta.get("underline"):
                rPr.set('u', '1')
            if run_meta.get("superscript"):
                rPr.set('baseline', '30000')
            if run_meta.get("subscript"):
                rPr.set('baseline', '-25000')

            font_name = run_meta.get("font_name")
            if font_name:
                latin = etree.SubElement(rPr, qn('a:latin'))
                latin.set('typeface', font_name)

            font_color = run_meta.get("font_color")
            if font_color:
                if font_color.startswith("theme:"):
                    solid = etree.SubElement(rPr, qn('a:solidFill'))
                    sc = etree.SubElement(solid, qn('a:schemeClr'))
                    sc.set('val', font_color[6:])
                elif _is_hex_color(font_color):
                    solid = etree.SubElement(rPr, qn('a:solidFill'))
                    sc = etree.SubElement(solid, qn('a:srgbClr'))
                    sc.set('val', font_color.upper())
            else:
                # Default to dk1
                solid = etree.SubElement(rPr, qn('a:solidFill'))
                sc = etree.SubElement(solid, qn('a:schemeClr'))
                sc.set('val', 'dk1')

            t = etree.SubElement(r, qn('a:t'))
            t.text = text


def _apply_body_props_xml(el, body_props):
    """Apply body properties (wrap, autofit, insets) to a group child."""
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    from pptx.oxml.ns import qn

    txBody = el.find(qn('p:txBody'))
    if txBody is None:
        return
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return

    for attr in ('wrap', 'vert', 'anchor', 'rtlCol'):
        val = body_props.get(attr)
        if val is not None:
            bodyPr.set(attr, str(val))

    for attr in ('lIns', 'rIns', 'tIns', 'bIns'):
        val = body_props.get(attr)
        if val is not None:
            bodyPr.set(attr, str(val))

    autofit = body_props.get('autofit')
    if autofit == 'spAutoFit':
        for child in list(bodyPr):
            if 'Autofit' in child.tag:
                bodyPr.remove(child)
        etree.SubElement(bodyPr, qn('a:spAutoFit'))
    elif autofit == 'normAutofit':
        for child in list(bodyPr):
            if 'Autofit' in child.tag:
                bodyPr.remove(child)
        norm = etree.SubElement(bodyPr, qn('a:normAutofit'))
        fs = body_props.get('fontScale')
        if fs: norm.set('fontScale', str(fs))
        ls = body_props.get('lnSpcReduction')
        if ls: norm.set('lnSpcReduction', str(ls))


def _parse_type_value(type_str):
    """Extract integer value from type string like 'AUTO_SHAPE (1)'."""
    if not type_str:
        return None
    parsed = re.match(r'.*?\((\d+)\)', type_str)
    if parsed:
        return int(parsed.group(1))
    return None


def _add_text_box_shape(slide, meta, images_dir, x, y, w, h, rotation):
    """Add a text box shape (no fill, no line by default)."""
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.rotation = rotation
    _remove_shape_style(shape)

    fill_meta = meta.get("fill")
    if fill_meta:
        _apply_fill(shape, _with_images_dir(fill_meta, images_dir))
    else:
        # Text boxes have no fill by default
        shape.fill.background()

    line_meta = meta.get("line")
    if line_meta:
        _apply_line(shape, line_meta)
    else:
        try:
            shape.line.fill.background()
        except Exception:
            pass

    text_meta = meta.get("text")
    if text_meta:
        _apply_text(shape, text_meta)

    # Apply body properties (wrap, autofit)
    body_props = meta.get("body_props")
    if body_props:
        _apply_body_props(shape, body_props)


def _add_auto_shape(slide, meta, images_dir, x, y, w, h, rotation):
    """Add an auto shape (rectangle, rounded rect, etc.) with proper fill/line."""
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    auto_type_str = meta.get("auto_shape_type", "")
    auto_type_val = None
    if auto_type_str:
        parsed = re.match(r'.*?\((\d+)\)', auto_type_str)
        if parsed:
            auto_type_val = int(parsed.group(1))

    try:
        auto_type = MSO_AUTO_SHAPE_TYPE(auto_type_val)
    except Exception:
        auto_type = MSO_AUTO_SHAPE_TYPE.RECTANGLE
    try:
        shape = slide.shapes.add_shape(auto_type, x, y, w, h)
    except Exception:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)

    shape.rotation = rotation
    _remove_shape_style(shape)

    # Apply fill
    fill_meta = meta.get("fill")
    if fill_meta:
        _apply_fill(shape, _with_images_dir(fill_meta, images_dir))
    else:
        # No fill metadata → transparent
        try:
            shape.fill.background()
        except Exception:
            pass

    # Apply line
    line_meta = meta.get("line")
    if line_meta:
        _apply_line(shape, line_meta)
    else:
        try:
            shape.line.fill.background()
        except Exception:
            pass

    # Apply text
    text_meta = meta.get("text")
    if text_meta:
        _apply_text(shape, text_meta)

    # Apply body properties (wrap, autofit, insets)
    body_props = meta.get("body_props")
    if body_props:
        _apply_body_props(shape, body_props)


def _with_images_dir(fill_meta, images_dir):
    """Return fill metadata augmented with the image directory for blip fills."""
    if not fill_meta or images_dir is None:
        return fill_meta
    result = dict(fill_meta)
    result["_images_dir"] = str(images_dir)
    return result


def _apply_raw_fill_to_spPr(spPr, raw_xml):
    """Replace fill XML under spPr with the original DrawingML fill element."""
    try:
        fill_el = etree.fromstring(raw_xml.encode("utf-8"))
    except Exception:
        return
    fill_names = {"solidFill", "noFill", "gradFill", "pattFill", "blipFill"}
    for child in list(spPr):
        local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if local in fill_names:
            spPr.remove(child)

    insert_at = len(spPr)
    for idx, child in enumerate(spPr):
        local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if local == "ln":
            insert_at = idx
            break
        if local == "prstGeom":
            insert_at = idx + 1
    spPr.insert(insert_at, fill_el)


def _apply_raw_line_to_spPr(spPr, raw_xml):
    """Replace line XML under spPr with the original DrawingML line element."""
    try:
        line_el = etree.fromstring(raw_xml.encode("utf-8"))
    except Exception:
        return
    for child in list(spPr):
        local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if local == "ln":
            spPr.remove(child)
    spPr.append(line_el)


def _add_image_shape(slide, meta, images_dir, x, y, w, h):
    """Add an image shape."""
    image_meta = meta.get("image", {})
    content_type = image_meta.get("content_type", "")
    filename = image_meta.get("filename", "")

    # Use stored filename if available
    if filename and images_dir.exists():
        img_path = Path(str(images_dir)) / filename
        if img_path.exists():
            try:
                slide.shapes.add_picture(str(img_path), x, y, w, h)
                return
            except Exception:
                pass

    # Fallback: search by extension
    ext = ".png"
    if "jpeg" in content_type:
        ext = ".jpg"
    elif "gif" in content_type:
        ext = ".gif"

    if images_dir.exists():
        for img_file in images_dir.iterdir():
            if img_file.suffix.lower() in [ext, ".png", ".jpg", ".jpeg"]:
                try:
                    slide.shapes.add_picture(str(img_file), x, y, w, h)
                    return
                except Exception:
                    continue

    # If no image found, add a placeholder rectangle
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.background()
    shape.text_frame.text = "[Image: {}]".format(content_type)


def _add_line_shape(slide, meta, x, y, w, h, rotation):
    """Add a line/connector shape with arrowheads."""
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    from pptx.oxml.ns import qn
    from lxml import etree

    # Determine connector type from metadata
    line_meta = meta.get("line", {})
    connector_geom = line_meta.get("connector_geom", "")

    conn_type = MSO_CONNECTOR_TYPE.STRAIGHT
    if connector_geom and connector_geom.startswith("bentConnector"):
        conn_type = MSO_CONNECTOR_TYPE.ELBOW

    try:
        shape = slide.shapes.add_connector(conn_type, x, y, x + w, y + h)
        _remove_shape_style(shape)
    except Exception:
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, max(w, 1), max(h, 1))
        shape.fill.background()

    # Apply flip flags on xfrm
    try:
        spPr = _get_spPr(shape)
        if spPr is not None:
            xfrm = spPr.find('{%s}xfrm' % A_NS)
            if xfrm is not None:
                for attr in ('flipH', 'flipV'):
                    val = line_meta.get(attr)
                    if val:
                        xfrm.set(attr, val)
    except Exception:
        pass

    # Apply line properties including arrowheads
    if line_meta:
        _apply_line(shape, line_meta)
        _apply_arrowheads(shape, line_meta)


def _apply_arrowheads(shape, line_meta):
    """Apply arrowhead (tailEnd/headEnd) to a connector shape's line XML."""
    try:
        spPr = _get_spPr(shape)
        if spPr is None:
            return
        ln = _make_ln(spPr, line_meta.get("width", 0))

        for end_key in ["tailEnd", "headEnd"]:
            arrow_info = line_meta.get(end_key)
            if not arrow_info:
                continue
            arr_type = arrow_info.get("type")
            if not arr_type:
                continue
            tag = "{{{}}}{}".format(A_NS, end_key)
            existing = ln.find(tag)
            if existing is not None:
                ln.remove(existing)
            elem = ln.makeelement(tag, {})
            elem.set("type", arr_type)
            if arrow_info.get("w"):
                elem.set("w", arrow_info["w"])
            if arrow_info.get("len"):
                elem.set("len", arrow_info["len"])
            ln.append(elem)
    except Exception:
        pass


def _add_table_shape(slide, meta, x, y, w, h):
    """Add a table shape."""
    table_meta = meta.get("table", {})
    rows = table_meta.get("rows", 1)
    cols = table_meta.get("cols", 1)

    slide.shapes.add_table(rows, cols, x, y, w, h)


def _apply_fill(shape, fill_meta):
    """Apply fill to a shape.

    Handles solid (RGB), scheme (theme), and no-fill types.
    """
    raw_xml = fill_meta.get("xml")
    if raw_xml:
        spPr = _get_spPr(shape)
        if spPr is not None:
            _apply_raw_fill_to_spPr(spPr, raw_xml)
            return

    fill_type = fill_meta.get("type")
    color = fill_meta.get("color")

    if fill_type == "none":
        try:
            shape.fill.background()
        except Exception:
            pass
        return

    if fill_type == "solid":
        if color:
            try:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(color)
            except Exception:
                pass
        return

    if fill_type == "scheme":
        # Use resolved absolute color if available (avoids theme dependency issues)
        resolved = fill_meta.get("_resolved")
        if resolved and _is_hex_color(resolved):
            try:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(resolved)
                return
            except Exception:
                pass
        _apply_scheme_fill(shape, color, fill_meta.get("modifiers"))
        return

    if fill_type == "blip":
        _apply_picture_fill(shape, fill_meta)
        return


def _apply_picture_fill(shape, fill_meta):
    """Apply an image as a shape fill using DrawingML blipFill."""
    filename = fill_meta.get("filename")
    if not filename:
        return

    image_path = Path(filename)
    if not image_path.is_absolute():
        images_dir = fill_meta.get("_images_dir")
        if not images_dir:
            return
        image_path = Path(images_dir) / filename
    if not image_path.exists():
        return

    try:
        _, r_id = shape.part.get_or_add_image_part(str(image_path))

        spPr = _get_spPr(shape)
        if spPr is None:
            return
        for child in list(spPr):
            local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if local in ("solidFill", "noFill", "gradFill", "pattFill", "blipFill"):
                spPr.remove(child)

        blip_fill = etree.Element("{{{}}}blipFill".format(A_NS))
        blip = etree.SubElement(blip_fill, "{{{}}}blip".format(A_NS))
        blip.set(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
            r_id,
        )
        stretch = etree.SubElement(blip_fill, "{{{}}}stretch".format(A_NS))
        etree.SubElement(stretch, "{{{}}}fillRect".format(A_NS))

        insert_at = 1 if len(spPr) > 1 else len(spPr)
        spPr.insert(insert_at, blip_fill)
    except Exception:
        pass


def _apply_scheme_fill(shape, theme_color, modifiers=None):
    """Apply a theme/scheme color fill via XML.

    Example: <a:solidFill><a:schemeClr val="accent6"/></a:solidFill>
    """
    if not theme_color:
        return
    try:
        spPr = _get_spPr(shape)
        if spPr is None:
            return
        # Remove existing fill elements
        fill_tags = ['{}Fill'.format(A_NS), '{}fill'.format(A_NS),
                     '{}NoFill'.format(A_NS), '{}noFill'.format(A_NS),
                     '{}solidFill'.format(A_NS), '{}blipFill'.format(A_NS),
                     '{}gradFill'.format(A_NS), '{}pattFill'.format(A_NS)]
        for child in list(spPr):
            if child.tag in fill_tags:
                spPr.remove(child)
        # Add scheme fill as first child of spPr (after xfrm, before geom)
        solid = etree.Element('{{{}}}solidFill'.format(A_NS))
        scheme = etree.SubElement(solid, '{{{}}}schemeClr'.format(A_NS))
        scheme.set('val', theme_color)
        if modifiers:
            for mod_name, mod_val in modifiers.items():
                mod = etree.SubElement(scheme, '{{{}}}{}'.format(A_NS, mod_name))
                mod.set('val', mod_val)
        spPr.insert(1, solid)  # After xfrm (index 0)
    except Exception:
        pass


def _is_hex_color(color_str):
    """Check if a color string is a valid hex color (6 hex digits)."""
    if not color_str or not isinstance(color_str, str):
        return False
    return bool(re.match(r'^[0-9A-Fa-f]{6}$', color_str))


def _apply_line(shape, line_meta):
    """Apply line properties to a shape.

    Handles connectors and regular shapes. Creates a:ln XML element
    if needed for arrowhead support.
    """
    raw_xml = line_meta.get("xml")
    if raw_xml:
        spPr = _get_spPr(shape)
        if spPr is not None:
            _apply_raw_line_to_spPr(spPr, raw_xml)
            return

    width = line_meta.get("width", 0)
    color = line_meta.get("color")

    if color is None:
        # No line
        try:
            shape.line.fill.background()
        except Exception:
            pass
        return

    if width > 0:
        try:
            shape.line.width = width
        except Exception:
            pass

    if _is_hex_color(color):
        try:
            shape.line.color.rgb = RGBColor.from_string(color)
        except Exception:
            pass
    else:
        # Theme color (e.g. "tx1") - set via XML
        _apply_theme_line_color(shape, color, width)


def _apply_theme_line_color(shape, theme_color, width, modifiers=None):
    """Apply a theme/scheme color to line via XML."""
    try:
        spPr = _get_spPr(shape)
        if spPr is None:
            return
        ln = spPr.find('{{{}}}ln'.format(A_NS))
        if ln is None:
            ln = _make_ln(spPr, width)
        # Remove existing fill
        for child in list(ln):
            tag = child.tag
            if 'Fill' in tag or 'fill' in tag or 'NoFill' in tag:
                ln.remove(child)
        solid = ln.makeelement('{{{}}}solidFill'.format(A_NS), {})
        scheme = solid.makeelement('{{{}}}schemeClr'.format(A_NS), {'val': theme_color})
        if modifiers:
            for mod_name, mod_val in modifiers.items():
                mod = scheme.makeelement('{{{}}}{}'.format(A_NS, mod_name), {'val': mod_val})
                scheme.append(mod)
        solid.append(scheme)
        ln.append(solid)
    except Exception:
        pass


def _get_spPr(shape):
    """Get the spPr element from a shape (handles both sp and cxnSp)."""
    return shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')


def _make_ln(spPr, width=0):
    """Create an a:ln element under spPr if it doesn't exist."""
    ln = spPr.find('{{{}}}ln'.format(A_NS))
    if ln is not None:
        return ln
    ln = etree.SubElement(spPr, '{{{}}}ln'.format(A_NS))
    if width > 0:
        ln.set('w', str(width))
    return ln


def _remove_shape_style(shape):
    """Remove the style element from a shape.

    Office PowerPoint misresolves schemeClr colors (tx1, dk1, etc.)
    as white when a <p:style> element is present. Removing it fixes
    text color rendering while also removing default shadows.
    """
    try:
        style = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}style')
        if style is not None:
            shape._element.remove(style)
    except Exception:
        pass


def _remove_theme_object_defaults(prs):
    """Remove <a:objectDefaults> from the theme.

    python-pptx's default theme includes <a:objectDefaults><a:spDef><a:style>
    with fontRef idx="minor" schemeClr="lt1". This interferes with Office's
    resolution of schemeClr colors (same mechanism as per-shape <p:style>).
    """
    try:
        A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        for part in prs.part.package.iter_parts():
            if '/ppt/theme/' in str(part.partname):
                theme_xml = etree.fromstring(part.blob)
                od = theme_xml.find('{%s}objectDefaults' % A_NS)
                if od is not None:
                    theme_xml.remove(od)
                    part.blob = etree.tostring(
                        theme_xml, xml_declaration=True, encoding='UTF-8', standalone=True
                    )
    except Exception:
        pass


def _ensure_line_xml(shape):
    """Ensure a minimal a:ln element exists for arrowhead support."""
    try:
        spPr = _get_spPr(shape)
        if spPr is None:
            return
        ln = spPr.find('{{{}}}ln'.format(A_NS))
        if ln is None:
            etree.SubElement(spPr, '{{{}}}ln'.format(A_NS))
    except Exception:
        pass


def _apply_run_font_color(run, color_value):
    """Apply font color to a run. Supports 'rgb:XXXXXX' and 'theme:XXX' formats."""
    if not color_value:
        return
    if color_value.startswith("theme:"):
        theme_val = color_value[6:]
        try:
            rPr_tag = '{{{}}}rPr'.format(A_NS)
            rPr = run._r.find(rPr_tag)
            if rPr is None:
                # Insert rPr BEFORE a:t (SubElement appends, which is invalid XML)
                rPr = etree.Element(rPr_tag)
                run._r.insert(0, rPr)
            # Remove existing fill
            for child in list(rPr):
                if 'Fill' in child.tag or 'fill' in child.tag:
                    rPr.remove(child)
            # Insert solidFill at position 0 (before latin/ea/sym etc.)
            # OOXML schema requires fill elements before font elements in rPr
            solid = etree.Element('{{{}}}solidFill'.format(A_NS))
            rPr.insert(0, solid)
            scheme = etree.SubElement(solid, '{{{}}}schemeClr'.format(A_NS))
            scheme.set('val', theme_val)
        except Exception:
            pass
    else:
        try:
            run.font.color.rgb = RGBColor.from_string(color_value)
        except Exception:
            pass


def _apply_body_props(shape, body_props):
    """Apply body properties (wrap, autofit, insets) to a shape."""
    try:
        A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        # Try both A_NS and P_NS for txBody
        txBody = shape._element.find('.//{%s}txBody' % A_NS)
        if txBody is None:
            txBody = shape._element.find('.//{%s}txBody' % P_NS)
        if txBody is None:
            return
        bodyPr = txBody.find('{%s}bodyPr' % A_NS)
        if bodyPr is None:
            return

        for attr in ('wrap', 'vert', 'anchor', 'rtlCol'):
            val = body_props.get(attr)
            if val is not None:
                bodyPr.set(attr, str(val))

        # Apply inset attributes (margins)
        for attr in ('lIns', 'rIns', 'tIns', 'bIns'):
            val = body_props.get(attr)
            if val is not None:
                bodyPr.set(attr, str(val))

        # Apply autofit
        autofit = body_props.get('autofit')
        if autofit == 'spAutoFit':
            # Remove any existing autofit elements
            for child in list(bodyPr):
                if 'Autofit' in child.tag or 'autofit' in child.tag.lower():
                    bodyPr.remove(child)
            spAutoFit = etree.SubElement(bodyPr, '{%s}spAutoFit' % A_NS)
        elif autofit == 'normAutofit':
            for child in list(bodyPr):
                if 'Autofit' in child.tag or 'autofit' in child.tag.lower():
                    bodyPr.remove(child)
            normAutofit = etree.SubElement(bodyPr, '{%s}normAutofit' % A_NS)
            fontScale = body_props.get('fontScale')
            lnSpcReduction = body_props.get('lnSpcReduction')
            if fontScale is not None:
                normAutofit.set('fontScale', str(fontScale))
            if lnSpcReduction is not None:
                normAutofit.set('lnSpcReduction', str(lnSpcReduction))
        elif autofit == 'noAutofit':
            for child in list(bodyPr):
                if 'Autofit' in child.tag or 'autofit' in child.tag.lower():
                    bodyPr.remove(child)
            etree.SubElement(bodyPr, '{%s}noAutofit' % A_NS)
    except Exception:
        pass


def _apply_text(shape, text_meta):
    """Apply text content and formatting to a shape."""
    paragraphs = text_meta.get("paragraphs", [])
    if not paragraphs:
        return

    tf = shape.text_frame
    if tf.paragraphs:
        tf.paragraphs[0].text = ""

    for i, para_meta in enumerate(paragraphs):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        level = para_meta.get("level", 0)
        para.level = level

        alignment = para_meta.get("alignment")
        from pptx.enum.text import PP_ALIGN
        if alignment:
            try:
                align_map = {
                    "LEFT (1)": PP_ALIGN.LEFT,
                    "CENTER (2)": PP_ALIGN.CENTER,
                    "RIGHT (3)": PP_ALIGN.RIGHT,
                    "JUSTIFY (4)": PP_ALIGN.JUSTIFY,
                }
                para.alignment = align_map.get(alignment, PP_ALIGN.LEFT)
            except Exception:
                pass
        else:
            # Override python-pptx default (CENTER for add_shape) to LEFT
            try:
                para.alignment = PP_ALIGN.LEFT
            except Exception:
                pass

        for run_meta in para_meta.get("runs", []):
            text = run_meta.get("text", "")
            run = para.add_run()
            run.text = text

            if run_meta.get("font_size"):
                run.font.size = run_meta["font_size"]
            if run_meta.get("bold"):
                run.font.bold = True
            if run_meta.get("italic"):
                run.font.italic = True
            if run_meta.get("underline"):
                run.font.underline = True
            # Superscript/subscript via baseline
            if run_meta.get("superscript"):
                rPr = run._r.find('{%s}rPr' % A_NS)
                if rPr is None:
                    rPr = etree.Element('{%s}rPr' % A_NS)
                    run._r.insert(0, rPr)
                rPr.set('baseline', '30000')
            if run_meta.get("subscript"):
                rPr = run._r.find('{%s}rPr' % A_NS)
                if rPr is None:
                    rPr = etree.Element('{%s}rPr' % A_NS)
                    run._r.insert(0, rPr)
                rPr.set('baseline', '-25000')
            if run_meta.get("font_name"):
                run.font.name = run_meta["font_name"]
            if run_meta.get("font_color"):
                _apply_run_font_color(run, run_meta["font_color"])
            # No explicit font_color: let theme inheritance handle it
