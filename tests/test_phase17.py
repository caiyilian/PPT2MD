"""Tests for Phase 17: Chart Image Export."""

import sys
from unittest.mock import MagicMock, patch

from ppt2md.parser.chart_export import export_chart_as_image, export_charts_from_slide, COM_AVAILABLE


def test_com_available_on_windows():
    if sys.platform == "win32":
        # Just check the flag is set correctly
        assert isinstance(COM_AVAILABLE, bool)


def test_export_chart_no_com():
    with patch("ppt2md.parser.chart_export.COM_AVAILABLE", False):
        chart = MagicMock()
        result = export_chart_as_image(chart, "test.png")
        assert result is False


def test_export_charts_no_chart(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    results = export_charts_from_slide(prs2.slides[0], str(tmp_path), 1)
    assert len(results) == 0
