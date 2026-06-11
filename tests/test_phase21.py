"""Tests for Phase 21: Document Properties Frontmatter."""

from pptx import Presentation

from ppt2md.converter.frontmatter import extract_document_properties, generate_frontmatter


def test_extract_properties(tmp_path):
    prs = Presentation()
    prs.core_properties.title = "Test Presentation"
    prs.core_properties.author = "Test Author"
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    props = extract_document_properties(prs2)
    assert props["title"] == "Test Presentation"
    assert props["author"] == "Test Author"


def test_extract_properties_empty():
    prs = Presentation()
    props = extract_document_properties(prs)
    assert props["title"] == ""
    assert props["author"] == ""


def test_generate_frontmatter(tmp_path):
    prs = Presentation()
    prs.core_properties.title = "My Talk"
    prs.core_properties.author = "Alice"
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    fm = generate_frontmatter(prs2, "test.pptx")
    assert fm.startswith("---")
    assert fm.endswith("---")
    assert "title: \"My Talk\"" in fm
    assert "author: \"Alice\"" in fm
    assert "source: \"test.pptx\"" in fm


def test_generate_frontmatter_no_title():
    prs = Presentation()
    fm = generate_frontmatter(prs)
    assert "---" in fm
    assert "title" not in fm
