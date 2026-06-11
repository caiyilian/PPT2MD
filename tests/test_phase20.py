"""Tests for Phase 20: Background Image."""

from pptx import Presentation

from ppt2md.parser.background import extract_background_image, format_background_markdown


def test_extract_background_no_bg(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    result = extract_background_image(prs2.slides[0], str(tmp_path), 1)
    assert result is None


def test_format_background_markdown_none():
    assert format_background_markdown(None) == ""


def test_format_background_markdown():
    info = {"filename": "slide_01_bg.png"}
    result = format_background_markdown(info)
    assert "background" in result
    assert "slide_01_bg.png" in result
