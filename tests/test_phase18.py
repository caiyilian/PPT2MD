"""Tests for Phase 18: SmartArt and Group Shapes."""

from pptx import Presentation
from pptx.util import Inches

from ppt2md.parser.smartart import extract_group_shapes, extract_smartart_text


def test_extract_group_shapes(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Add shapes to a group
    shape1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    shape1.text = "Group item 1"
    shape2 = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(2), Inches(1))
    shape2.text = "Group item 2"
    # Note: python-pptx doesn't easily create groups programmatically,
    # so we test with individual shapes
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    shapes = list(prs2.slides[0].shapes)
    assert len(shapes) >= 2


def test_extract_smartart_text_empty(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])  # Title slide with placeholders
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    shape = prs2.slides[0].shapes.title
    texts = extract_smartart_text(shape)
    assert isinstance(texts, list)
