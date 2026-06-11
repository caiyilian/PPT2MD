"""Tests for Phase 19: Media and OLE Objects."""

from pptx import Presentation

from ppt2md.parser.media import (
    extract_media_from_slide,
    detect_ole_objects,
    format_media_markdown,
    format_ole_markdown,
)


def test_extract_media_no_media(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    media = extract_media_from_slide(prs2.slides[0], str(tmp_path / "media"), 1)
    assert len(media) == 0


def test_detect_ole_empty(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    ole = detect_ole_objects(prs2.slides[0])
    assert len(ole) == 0


def test_format_media_markdown_audio():
    info = {"filename": "test.mp3", "content_type": "audio/mpeg"}
    result = format_media_markdown(info)
    assert "Audio" in result
    assert "test.mp3" in result


def test_format_media_markdown_video():
    info = {"filename": "test.mp4", "content_type": "video/mp4"}
    result = format_media_markdown(info)
    assert "Video" in result


def test_format_ole_markdown():
    info = {"name": "Document1"}
    result = format_ole_markdown(info)
    assert "Embedded Object" in result
    assert "Document1" in result
