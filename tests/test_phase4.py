"""Tests for Phase 4: Text Formatting to Markdown."""

from pptx import Presentation
from pptx.util import Pt

from ppt2md.converter.format_utils import apply_inline_formatting, format_paragraph_runs


def test_bold():
    assert apply_inline_formatting("hello", bold=True) == "**hello**"


def test_italic():
    assert apply_inline_formatting("hello", italic=True) == "*hello*"


def test_underline():
    assert apply_inline_formatting("hello", underline=True) == "<u>hello</u>"


def test_strikethrough():
    assert apply_inline_formatting("hello", strikethrough=True) == "~~hello~~"


def test_bold_italic():
    result = apply_inline_formatting("hello", bold=True, italic=True)
    assert "**" in result
    assert "*" in result
    assert "hello" in result


def test_no_formatting():
    assert apply_inline_formatting("hello") == "hello"


def test_empty_text():
    assert apply_inline_formatting("") == ""
    assert apply_inline_formatting("", bold=True) == ""


def test_format_paragraph_runs_bold(tmp_path):
    pptx = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Test"
    para = title.text_frame.paragraphs[0]
    para.runs[0].font.bold = True
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    slide2 = prs2.slides[0]
    title2 = slide2.shapes.title
    result = format_paragraph_runs(title2.text_frame.paragraphs[0])
    assert "**Test**" == result


def test_format_paragraph_runs_mixed(tmp_path):
    pptx = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    # Clear and add multiple runs
    title.text = ""
    tf = title.text_frame
    para = tf.paragraphs[0]
    run1 = para.add_run()
    run1.text = "Bold "
    run1.font.bold = True
    run2 = para.add_run()
    run2.text = "Normal"
    run2.font.bold = False
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    result = format_paragraph_runs(prs2.slides[0].shapes.title.text_frame.paragraphs[0])
    assert "**Bold **Normal" == result
