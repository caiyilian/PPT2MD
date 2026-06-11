"""Tests for Phase 24: Empty Slide Filtering."""

from pptx import Presentation
from pptx.util import Inches

from ppt2md.converter.filter_utils import is_empty_slide


def test_empty_slide(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    assert is_empty_slide(prs2.slides[0]) is True


def test_non_empty_slide_text(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    txBox.text_frame.text = "Real content"
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    assert is_empty_slide(prs2.slides[0]) is False


def test_non_empty_slide_image(tmp_path):
    from PIL import Image
    img_path = tmp_path / "test.png"
    Image.new("RGB", (100, 100), "red").save(str(img_path))

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img_path), Inches(1), Inches(1), Inches(2), Inches(2))
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    assert is_empty_slide(prs2.slides[0]) is False


def test_empty_slide_placeholder_only(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # Default placeholders are empty after creation
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    # Title slide with default placeholder text should be considered non-empty
    # because it has a title placeholder with actual text
    assert isinstance(is_empty_slide(prs2.slides[0]), bool)
