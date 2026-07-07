"""Tests for PPTX to HTML conversion."""

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.util import Centipoints, Inches, Pt

from ppt2md.converter.html import (
    HtmlRenderContext,
    convert_pptx_to_html,
    main as ppt2html_main,
    _render_group,
    _stroke_info,
)


def _create_html_fixture_pptx(tmp_path):
    image_path = tmp_path / "sample.png"
    Image.new("RGB", (80, 60), "navy").save(image_path)

    prs = Presentation()
    prs.slide_width = Inches(8)
    prs.slide_height = Inches(4.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 246, 250)

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.4),
        Inches(0.35),
        Inches(2.5),
        Inches(0.8),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(248, 203, 173)
    box.line.color.rgb = RGBColor(197, 90, 17)
    box.line.width = Pt(1.5)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = ""
    run = paragraph.add_run()
    run.text = "HTML Export"
    run.font.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(32, 55, 100)

    ellipse = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(3.25),
        Inches(0.35),
        Inches(0.9),
        Inches(0.9),
    )
    ellipse.fill.background()
    ellipse.line.color.rgb = RGBColor(112, 48, 160)
    ellipse.line.width = Pt(2)

    slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(2.95),
        Inches(0.75),
        Inches(3.25),
        Inches(0.75),
    )

    slide.shapes.add_picture(
        str(image_path),
        Inches(0.5),
        Inches(1.6),
        Inches(1.4),
        Inches(1.0),
    )

    table_shape = slide.shapes.add_table(
        2,
        2,
        Inches(2.3),
        Inches(1.55),
        Inches(2.6),
        Inches(1.25),
    )
    table = table_shape.table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Accuracy"
    table.cell(1, 1).text = "98%"

    formula_box = slide.shapes.add_textbox(
        Inches(5.2),
        Inches(1.55),
        Inches(2.1),
        Inches(0.6),
    )
    formula_box.text_frame.text = "$E=mc^2$"

    pptx_path = tmp_path / "fixture.pptx"
    prs.save(pptx_path)
    return pptx_path


def test_convert_pptx_to_html_embeds_core_content(tmp_path):
    pptx_path = _create_html_fixture_pptx(tmp_path)
    output_path = tmp_path / "fixture.html"

    result = convert_pptx_to_html(pptx_path, output_path)

    html = output_path.read_text(encoding="utf-8")
    assert result == output_path
    assert "HTML Export" in html
    assert "Accuracy" in html
    assert "$E=mc^2$" in html
    assert "data:image/png;base64," in html
    assert "https://cdn.jsdelivr.net/npm/katex" in html
    assert "class=\"ppt-shape\"" in html
    assert "class=\"ppt-line\"" in html
    assert "left:" in html and "top:" in html


def test_ppt2html_cli_writes_requested_output(tmp_path):
    pptx_path = _create_html_fixture_pptx(tmp_path)
    output_path = tmp_path / "cli-output.html"

    exit_code = ppt2html_main([str(pptx_path), "-o", str(output_path)])

    assert exit_code == 0
    assert output_path.exists()
    assert "cli-output" not in output_path.read_text(encoding="utf-8")


def test_render_group_preserves_centipoint_font_size():
    meta = {
        "type": "GROUP (6)",
        "position": {"x": 0, "y": 0},
        "size": {"width": 914400, "height": 914400},
        "group": {
            "coord_space": {
                "chOffX": 0,
                "chOffY": 0,
                "chExtCX": 914400,
                "chExtCY": 914400,
            },
            "children": [
                {
                    "name": "group text",
                    "type": "AUTO_SHAPE (1)",
                    "auto_shape_type": "ROUNDED_RECTANGLE (5)",
                    "position": {"x": 0, "y": 0},
                    "size": {"width": 457200, "height": 228600},
                    "body_props": {"vert": "horz"},
                    "fill": {"type": "none"},
                    "text": {
                        "paragraphs": [
                            {
                                "runs": [
                                    {
                                        "text": "HFEM",
                                        "font_size": Centipoints(1400),
                                    }
                                ]
                            }
                        ]
                    },
                }
            ],
        },
    }
    ctx = HtmlRenderContext({}, {})

    html = "\n".join(_render_group(meta, ctx, 1))

    assert "font-size:14.000pt" in html
    assert "font-size:1778.000pt" not in html
    assert "writing-mode:vertical-rl" not in html


def test_render_text_handles_vert270_without_writing_mode():
    meta = {
        "name": "rotated text",
        "type": "AUTO_SHAPE (1)",
        "auto_shape_type": "ROUNDED_RECTANGLE (5)",
        "position": {"x": 0, "y": 0},
        "size": {"width": 457200, "height": 228600},
        "body_props": {"vert": "vert270"},
        "text": {
            "paragraphs": [
                {
                    "runs": [
                        {
                            "text": "HFEM",
                            "font_size": Centipoints(1400),
                        }
                    ]
                }
            ]
        },
    }
    ctx = HtmlRenderContext({}, {})

    html = "\n".join(_render_group({
        "type": "GROUP (6)",
        "position": {"x": 0, "y": 0},
        "size": {"width": 914400, "height": 914400},
        "group": {
            "coord_space": {
                "chOffX": 0,
                "chOffY": 0,
                "chExtCX": 914400,
                "chExtCY": 914400,
            },
            "children": [meta],
        },
    }, ctx, 1))

    assert "transform:rotate(-90deg)" in html
    assert "writing-mode:vertical-rl" not in html


def test_stroke_info_applies_theme_luminance_modifier():
    line_meta = {
        "color": "bg1",
        "width": 19050,
        "dash": "sysDash",
        "xml": (
            '<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'w="19050"><a:solidFill><a:schemeClr val="bg1">'
            '<a:lumMod val="50000"/></a:schemeClr></a:solidFill>'
            '<a:prstDash val="sysDash"/></a:ln>'
        ),
    }

    stroke = _stroke_info(line_meta, {"lt1": "FFFFFF"})

    assert stroke["color"] == "#808080"
    assert stroke["dash"] == "6 3"
