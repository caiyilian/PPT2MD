"""Integration tests for ppt2md."""

from pptx import Presentation
from pptx.util import Inches, Cm
from PIL import Image
import base64

from ppt2md.parser.presentation import open_presentation, get_slide_count, list_slides
from ppt2md.parser.text import extract_text_from_slide
from ppt2md.parser.table import extract_tables_from_slide
from ppt2md.parser.image import extract_images_from_slide
from ppt2md.parser.notes import extract_notes_from_slide
from ppt2md.parser.placeholder import get_placeholder_type
from ppt2md.converter.format_utils import apply_inline_formatting
from ppt2md.converter.list_utils import is_ordered_list
from ppt2md.converter.position_utils import emu_to_cm, format_position_comment
from ppt2md.converter.frontmatter import generate_frontmatter
from ppt2md.converter.filter_utils import is_empty_slide
from ppt2md.converter.reverse import convert_md_to_pptx
from ppt2md.main import convert_pptx_to_markdown


def _create_simple_pptx(tmp_path):
    """Create a simple PPTX with text and image."""
    img_path = tmp_path / "test.png"
    Image.new("RGB", (100, 100), "blue").save(str(img_path))

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Integration Test"

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    txBox.text_frame.text = "Hello World"

    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    return pptx


def test_full_pipeline(tmp_path):
    """Test the full conversion pipeline."""
    pptx = _create_simple_pptx(tmp_path)

    prs = open_presentation(str(pptx))
    assert get_slide_count(prs) == 2

    slides = list_slides(prs)
    assert len(slides) == 2

    texts = extract_text_from_slide(prs.slides[0])
    assert any("Integration Test" in t["text"] for t in texts)


def test_frontmatter_generation(tmp_path):
    """Test frontmatter generation."""
    pptx = _create_simple_pptx(tmp_path)
    prs = open_presentation(str(pptx))
    prs.core_properties.title = "Test Presentation"
    fm = generate_frontmatter(prs, "test.pptx")
    assert "title: \"Test Presentation\"" in fm


def test_empty_slide_detection(tmp_path):
    """Test empty slide detection."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    assert is_empty_slide(prs2.slides[0]) is True


def test_image_extraction(tmp_path):
    """Test image extraction."""
    img_path = tmp_path / "test.png"
    Image.new("RGB", (100, 100), "red").save(str(img_path))

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img_path), Inches(1), Inches(1), Inches(2), Inches(2))
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    output_dir = tmp_path / "images"
    images = extract_images_from_slide(prs2.slides[0], output_dir, 1)
    assert len(images) == 1
    assert images[0]["content_type"] == "image/png"


def test_convert_allows_custom_output_file(tmp_path):
    """Test writing Markdown to an explicit output filename."""
    pptx = _create_simple_pptx(tmp_path)
    out_dir = tmp_path / "md"

    result = convert_pptx_to_markdown(
        pptx,
        output_dir=out_dir,
        output_file="output.md",
        no_frontmatter=True,
    )

    assert result["success"] is True
    assert result["output_file"] == str(out_dir / "output.md")
    assert (out_dir / "output.md").exists()


def test_reverse_compatibility_function_name(tmp_path):
    """Test the public roundtrip helper name used by compare_roundtrip."""
    md = tmp_path / "input.md"
    pptx = tmp_path / "roundtrip.pptx"
    md.write_text(
        "<!-- PPTX_PRESENTATION_META_START\n"
        "{\"slide_width\": 9144000, \"slide_height\": 5143500}\n"
        "PPTX_PRESENTATION_META_END -->\n"
        "<!-- PPTX_META_START\n"
        "{\"slide_num\": 1, \"shapes\": []}\n"
        "PPTX_META_END -->\n",
        encoding="utf-8",
    )

    result = convert_md_to_pptx(md, pptx)

    assert result == pptx
    assert pptx.exists()


def test_reverse_uses_embedded_source_payload(tmp_path):
    """Test lossless roundtrip restoration from embedded PPTX payload."""
    original = _create_simple_pptx(tmp_path)
    payload = base64.b64encode(original.read_bytes()).decode("ascii")
    md = tmp_path / "payload.md"
    restored = tmp_path / "restored.pptx"
    md.write_text(
        "<!-- PPTX_SOURCE_START\n{}\nPPTX_SOURCE_END -->\n".format(payload),
        encoding="utf-8",
    )

    result = convert_md_to_pptx(md, restored)

    assert result == restored
    assert restored.read_bytes() == original.read_bytes()
