"""Tests for Phase 3: Text Frame Extraction."""

from pptx import Presentation
from pptx.util import Inches

from ppt2md.parser.text import extract_text_from_slide, extract_text_from_shape


def _create_text_pptx(path):
    """Create a PPTX with text content."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # Title placeholder
    title = slide.shapes.title
    title.text = "Hello World"
    # Body placeholder
    body = slide.placeholders[1]
    body.text = "First paragraph\nSecond paragraph"
    prs.save(str(path))


def _create_empty_slide_pptx(path):
    """Create a PPTX with an empty slide."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    # Clear all text from placeholders
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                para.clear()
    prs.save(str(path))


def test_extract_text_from_slide(tmp_path):
    pptx = tmp_path / "test.pptx"
    _create_text_pptx(pptx)
    prs = Presentation(str(pptx))
    results = extract_text_from_slide(prs.slides[0])
    texts = [r["text"] for r in results]
    assert "Hello World" in texts


def test_extract_text_from_slide_paragraphs(tmp_path):
    pptx = tmp_path / "test.pptx"
    _create_text_pptx(pptx)
    prs = Presentation(str(pptx))
    results = extract_text_from_slide(prs.slides[0])
    texts = [r["text"] for r in results]
    assert "First paragraph" in texts
    assert "Second paragraph" in texts


def test_extract_text_skips_empty(tmp_path):
    pptx = tmp_path / "test.pptx"
    _create_empty_slide_pptx(pptx)
    prs = Presentation(str(pptx))
    results = extract_text_from_slide(prs.slides[0])
    assert len(results) == 0


def test_extract_text_from_shape(tmp_path):
    pptx = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    tf = txBox.text_frame
    tf.text = "Text box content"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    all_texts = []
    for shape in prs2.slides[0].shapes:
        all_texts.extend(extract_text_from_shape(shape))
    assert "Text box content" in all_texts


def test_extract_text_shape_without_text_frame(tmp_path):
    pptx = tmp_path / "test.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    slide = prs2.slides[0]
    # Find a non-text shape or use the placeholder
    for shape in slide.shapes:
        if not shape.has_text_frame:
            assert extract_text_from_shape(shape) == []
            return
    # All shapes have text frames, test with empty one
    assert extract_text_from_shape(slide.shapes[0]) == []
