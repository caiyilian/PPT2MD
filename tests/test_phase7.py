"""Tests for Phase 7: Image Extraction."""

import os

from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from ppt2md.parser.image import get_image_extension, extract_images_from_slide


def _create_image(path, size=(100, 100), color="red"):
    """Create a test image."""
    img = Image.new("RGB", size, color)
    img.save(str(path))


def _create_pptx_with_image(pptx_path, image_path):
    """Create a PPTX containing an image."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), Inches(2), Inches(2))
    prs.save(str(pptx_path))


def test_get_image_extension_png():
    assert get_image_extension("image/png") == ".png"


def test_get_image_extension_jpg():
    assert get_image_extension("image/jpeg") == ".jpg"


def test_get_image_extension_unknown():
    assert get_image_extension("image/unknown") == ".bin"


def test_extract_images(tmp_path):
    img_path = tmp_path / "test.png"
    _create_image(img_path)

    pptx_path = tmp_path / "test.pptx"
    _create_pptx_with_image(pptx_path, img_path)

    output_dir = tmp_path / "images"
    prs = Presentation(str(pptx_path))
    images = extract_images_from_slide(prs.slides[0], output_dir, 1)

    assert len(images) == 1
    assert images[0]["content_type"] == "image/png"
    assert os.path.exists(images[0]["path"])


def test_extract_images_filename_format(tmp_path):
    img_path = tmp_path / "test.png"
    _create_image(img_path)

    pptx_path = tmp_path / "test.pptx"
    _create_pptx_with_image(pptx_path, img_path)

    output_dir = tmp_path / "images"
    prs = Presentation(str(pptx_path))
    images = extract_images_from_slide(prs.slides[0], output_dir, 3)

    assert images[0]["filename"] == "slide_03_img_01.png"


def test_extract_images_empty_slide(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    prs.save(str(pptx_path))

    output_dir = tmp_path / "images"
    prs2 = Presentation(str(pptx_path))
    images = extract_images_from_slide(prs2.slides[0], output_dir, 1)

    assert len(images) == 0
