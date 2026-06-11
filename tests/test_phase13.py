"""Tests for Phase 13: Shape Position Sorting."""

from pptx import Presentation
from pptx.util import Cm
from PIL import Image

from ppt2md.converter.sort_utils import sort_shapes_by_position


def _create_pptx_with_positions(tmp_path):
    """Create a PPTX with shapes at different positions."""
    img_path = tmp_path / "test.png"
    Image.new("RGB", (50, 50), "red").save(str(img_path))

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Shape at bottom-left
    slide.shapes.add_picture(str(img_path), Cm(1), Cm(15), Cm(3), Cm(3))
    # Shape at top-right
    slide.shapes.add_picture(str(img_path), Cm(15), Cm(1), Cm(3), Cm(3))
    # Shape at top-left
    slide.shapes.add_picture(str(img_path), Cm(1), Cm(1), Cm(3), Cm(3))
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    return pptx


def test_sort_reading_order(tmp_path):
    pptx = _create_pptx_with_positions(tmp_path)
    prs = Presentation(str(pptx))
    shapes = list(prs.slides[0].shapes)
    sorted_shapes = sort_shapes_by_position(shapes)

    # Should be: top-left, top-right, bottom-left
    assert sorted_shapes[0].top < sorted_shapes[2].top
    assert sorted_shapes[0].left < sorted_shapes[1].left


def test_sort_empty():
    assert sort_shapes_by_position([]) == []


def test_sort_none_values():
    """Shapes with None positions should be sorted to end."""
    from unittest.mock import MagicMock
    s1 = MagicMock()
    s1.top = None
    s1.left = 100
    s2 = MagicMock()
    s2.top = 100
    s2.left = 100
    result = sort_shapes_by_position([s1, s2])
    assert result[0].top == 100
    assert result[1].top is None
