"""Tests for Phase 16: Chart Data Extraction."""

from pptx import Presentation

from ppt2md.parser.chart import extract_chart_as_markdown, extract_charts_from_slide


def test_extract_charts_no_chart(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    charts = extract_charts_from_slide(prs2.slides[0])
    assert len(charts) == 0


def test_extract_chart_as_markdown_no_data():
    from unittest.mock import MagicMock
    chart = MagicMock()
    chart.has_data = False
    result = extract_chart_as_markdown(chart)
    assert "no data" in result.lower()
