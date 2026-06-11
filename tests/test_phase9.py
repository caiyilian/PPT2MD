"""Tests for Phase 9: Image Crop Handling."""

from pptx import Presentation
from pptx.util import Cm
from PIL import Image

from ppt2md.converter.position_utils import get_crop_info, format_position_comment


def _create_pptx_with_image(tmp_path):
    img_path = tmp_path / "test.png"
    Image.new("RGB", (100, 100), "red").save(str(img_path))

    pptx = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img_path), Cm(5), Cm(10), Cm(15), Cm(8))
    prs.save(str(pptx))
    return pptx


def test_get_crop_info_no_crop(tmp_path):
    pptx = _create_pptx_with_image(tmp_path)
    prs = Presentation(str(pptx))
    shape = prs.slides[0].shapes[0]
    crop = get_crop_info(shape)
    # No crop set, should return None
    assert crop is None


def test_format_position_comment_without_crop(tmp_path):
    pptx = _create_pptx_with_image(tmp_path)
    prs = Presentation(str(pptx))
    shape = prs.slides[0].shapes[0]
    comment = format_position_comment(shape)
    assert "<!-- position:" in comment
    assert "cm" in comment
    assert "-->" in comment
    assert "crop" not in comment


def test_get_crop_info_non_image(tmp_path):
    pptx = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Title"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    shape = prs2.slides[0].shapes.title
    assert get_crop_info(shape) is None
