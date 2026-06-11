"""Tests for Phase 12: Speaker Notes."""

from pptx import Presentation

from ppt2md.parser.notes import extract_notes_from_slide, format_notes_markdown


def test_extract_notes(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These are speaker notes"
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    notes = extract_notes_from_slide(prs2.slides[0])
    assert notes == "These are speaker notes"


def test_extract_notes_empty(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    notes = extract_notes_from_slide(prs2.slides[0])
    assert notes == ""


def test_format_notes_markdown():
    result = format_notes_markdown("Some notes")
    assert "<!-- Notes -->" in result
    assert "### Speaker Notes" in result
    assert "Some notes" in result


def test_format_notes_markdown_empty():
    assert format_notes_markdown("") == ""
    assert format_notes_markdown(None) == ""
