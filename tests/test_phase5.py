"""Tests for Phase 5: Placeholder Type Mapping."""

from pptx import Presentation

from ppt2md.parser.placeholder import get_placeholder_type, placeholder_to_markdown_prefix


def test_title_placeholder(tmp_path):
    pptx = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    shape = prs2.slides[0].shapes.title
    assert get_placeholder_type(shape) == "title"


def test_body_placeholder(tmp_path):
    pptx = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section Header has BODY type
    body = slide.placeholders[1]
    body.text = "Body text"
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    shape = prs2.slides[0].placeholders[1]
    assert get_placeholder_type(shape) == "body"


def test_non_placeholder(tmp_path):
    pptx = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    from pptx.util import Inches
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    txBox.text = "Not a placeholder"
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    for shape in prs2.slides[0].shapes:
        if not shape.is_placeholder:
            assert get_placeholder_type(shape) == "object"
            break


def test_title_prefix():
    assert placeholder_to_markdown_prefix("title") == "## "


def test_subtitle_prefix():
    assert placeholder_to_markdown_prefix("subtitle") == "### "


def test_body_prefix():
    assert placeholder_to_markdown_prefix("body") == ""


def test_skip_prefix():
    assert placeholder_to_markdown_prefix("skip") == ""
