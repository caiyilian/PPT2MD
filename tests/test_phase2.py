"""Tests for Phase 2: Basic PPTX Parsing."""

from pptx import Presentation

from ppt2md.parser.presentation import open_presentation, get_slide_count, list_slides


def _create_test_pptx(path, slide_count=3):
    """Helper: create a simple PPTX with N blank slides."""
    prs = Presentation()
    for _ in range(slide_count):
        prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(path))


def test_open_presentation(tmp_path):
    pptx = tmp_path / "test.pptx"
    _create_test_pptx(pptx, 3)
    prs = open_presentation(str(pptx))
    assert prs is not None


def test_get_slide_count(tmp_path):
    pptx = tmp_path / "test.pptx"
    _create_test_pptx(pptx, 3)
    prs = open_presentation(str(pptx))
    assert get_slide_count(prs) == 3


def test_get_slide_count_single(tmp_path):
    pptx = tmp_path / "test.pptx"
    _create_test_pptx(pptx, 1)
    prs = open_presentation(str(pptx))
    assert get_slide_count(prs) == 1


def test_list_slides(tmp_path):
    pptx = tmp_path / "test.pptx"
    _create_test_pptx(pptx, 3)
    prs = open_presentation(str(pptx))
    slides = list_slides(prs)
    assert len(slides) == 3
    assert slides[0]["index"] == 0
    assert slides[1]["index"] == 1
    assert slides[2]["index"] == 2


def test_list_slides_shape_count(tmp_path):
    pptx = tmp_path / "test.pptx"
    _create_test_pptx(pptx, 2)
    prs = open_presentation(str(pptx))
    slides = list_slides(prs)
    for s in slides:
        assert "shape_count" in s
        assert s["shape_count"] >= 1  # slide_layouts[0] always has placeholders
