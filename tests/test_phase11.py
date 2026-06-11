"""Tests for Phase 11: Table Extraction."""

from pptx import Presentation
from pptx.util import Inches

from ppt2md.parser.table import extract_table_as_markdown, extract_tables_from_slide


def _create_pptx_with_table(tmp_path, rows=2, cols=3):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(4)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    for r in range(rows):
        for c in range(cols):
            table.cell(r, c).text = "R{}C{}".format(r, c)
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    return pptx


def test_extract_table_basic(tmp_path):
    pptx = _create_pptx_with_table(tmp_path)
    prs = Presentation(str(pptx))
    table = prs.slides[0].shapes[0].table
    md = extract_table_as_markdown(table)
    assert "| R0C0 | R0C1 | R0C2 |" in md
    assert "| R1C0 | R1C1 | R1C2 |" in md
    assert "| --- | --- | --- |" in md


def test_extract_table_single_row(tmp_path):
    pptx = _create_pptx_with_table(tmp_path, rows=1, cols=2)
    prs = Presentation(str(pptx))
    table = prs.slides[0].shapes[0].table
    md = extract_table_as_markdown(table)
    assert "| R0C0 | R0C1 |" in md
    assert "| --- | --- |" in md


def test_extract_tables_from_slide(tmp_path):
    pptx = _create_pptx_with_table(tmp_path)
    prs = Presentation(str(pptx))
    tables = extract_tables_from_slide(prs.slides[0])
    assert len(tables) == 1
    assert "R0C0" in tables[0]["markdown"]


def test_extract_tables_no_table(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    pptx = tmp_path / "test.pptx"
    prs.save(str(pptx))
    prs2 = Presentation(str(pptx))
    tables = extract_tables_from_slide(prs2.slides[0])
    assert len(tables) == 0
