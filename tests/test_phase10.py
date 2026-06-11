"""Tests for Phase 10: Image Deduplication."""

from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from ppt2md.parser.image import extract_images_from_slide


def _create_image(path, color="red"):
    Image.new("RGB", (100, 100), color).save(str(path))


def test_dedup_same_image(tmp_path):
    """Same image referenced on two slides should only be saved once."""
    img_path = tmp_path / "test.png"
    _create_image(img_path)

    pptx = tmp_path / "test.pptx"
    prs = Presentation()
    for _ in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img_path), Inches(1), Inches(1), Inches(2), Inches(2))
    prs.save(str(pptx))

    output_dir = tmp_path / "images"
    prs2 = Presentation(str(pptx))
    seen_rids = {}

    imgs1 = extract_images_from_slide(prs2.slides[0], output_dir, 1, seen_rids)
    imgs2 = extract_images_from_slide(prs2.slides[1], output_dir, 2, seen_rids)

    # First slide extracts the image
    assert len(imgs1) == 1
    assert imgs1[0]["deduplicated"] is False

    # Second slide reuses the same filename
    assert len(imgs2) == 1
    assert imgs2[0]["deduplicated"] is True
    assert imgs2[0]["filename"] == imgs1[0]["filename"]


def test_no_dedup_different_images(tmp_path):
    """Different images should not be deduplicated."""
    img1 = tmp_path / "red.png"
    img2 = tmp_path / "blue.png"
    _create_image(img1, "red")
    _create_image(img2, "blue")

    pptx = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img1), Inches(1), Inches(1), Inches(2), Inches(2))
    slide.shapes.add_picture(str(img2), Inches(4), Inches(1), Inches(2), Inches(2))
    prs.save(str(pptx))

    output_dir = tmp_path / "images"
    prs2 = Presentation(str(pptx))
    seen_rids = {}
    imgs = extract_images_from_slide(prs2.slides[0], output_dir, 1, seen_rids)

    assert len(imgs) == 2
    assert imgs[0]["deduplicated"] is False
    assert imgs[1]["deduplicated"] is False
    assert imgs[0]["filename"] != imgs[1]["filename"]
