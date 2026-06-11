"""Tests for Phase 8: Image Position and Size."""

from pptx import Presentation
from pptx.util import Inches, Cm
from PIL import Image

from ppt2md.converter.position_utils import emu_to_cm, emu_to_px, format_position_comment


def test_emu_to_cm():
    assert abs(emu_to_cm(914400) - 2.54) < 0.01


def test_emu_to_cm_none():
    assert emu_to_cm(None) == 0.0


def test_emu_to_px():
    assert abs(emu_to_px(914400, dpi=96) - 96.0) < 0.01


def test_emu_to_px_none():
    assert emu_to_px(None) == 0.0


def test_format_position_comment(tmp_path):
    img_path = tmp_path / "test.png"
    Image.new("RGB", (100, 100), "red").save(str(img_path))

    pptx = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img_path), Cm(5), Cm(10), Cm(15), Cm(8))
    prs.save(str(pptx))

    prs2 = Presentation(str(pptx))
    shape = prs2.slides[0].shapes[0]
    comment = format_position_comment(shape)
    assert "<!-- position:" in comment
    assert "cm" in comment
    assert "width=" in comment
    assert "height=" in comment
