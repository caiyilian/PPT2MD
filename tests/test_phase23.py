"""Tests for Phase 23: Master Layout Placeholder Filtering."""

from ppt2md.converter.filter_utils import is_default_placeholder, filter_shape_text, should_skip_shape


def test_is_default_placeholder_title():
    assert is_default_placeholder("click to add title") is True
    assert is_default_placeholder("Click to Add Title") is True


def test_is_default_placeholder_chinese():
    assert is_default_placeholder("单击此处添加标题") is True


def test_is_default_placeholder_custom():
    assert is_default_placeholder("My Custom Title") is False


def test_is_default_placeholder_empty():
    assert is_default_placeholder("") is False


def test_filter_shape_text_custom():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    txBox.text_frame.text = "Real content here"
    shape = txBox
    skip, text = filter_shape_text(shape)
    assert skip is False
    assert "Real content here" in text


def test_should_skip_shape_custom():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    txBox.text_frame.text = "Real content"
    assert should_skip_shape(txBox) is False
